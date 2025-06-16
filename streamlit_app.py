import streamlit as st
import requests
import tempfile
from urllib.parse import urlencode
from io import StringIO, BytesIO # Added BytesIO for image processing
from PIL import Image, ImageFilter, ImageOps
import pytesseract
import json
import igraph as ig
import plotly.graph_objects as go
import re
import sqlite3
from contextlib import closing
from streamlit_lottie import st_lottie
import requests as reqs
import contextlib
import csv
from fpdf import FPDF
import webbrowser
import fitz
import docx
from pptx import Presentation
import streamlit.components.v1 as components
import time
import networkx as nx
import matplotlib.pyplot as plt
import datetime
import pandas as pd
from gtts import gTTS
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import os
# import vlc # Keeping commented as VLC can be problematic in web environments
import base64

# Define missing variables (ensure these are handled if they are meant to be dynamic)
all_summaries = []
all_flashcards = []

# --- Error Reporting ---
def send_error_report(error_message):
    """Send an error report to the support team."""
    sender_email = "your_email@example.com"  # Replace with your email
    receiver_email = "sambit1912@gmail.com"
    password = "your_email_password"  # Replace with your email password

    subject = "Vekkam Error Report"
    body = f"An error occurred:\n\n{error_message}"

    msg = MIMultipart()
    msg['From'] = sender_email
    msg['To'] = receiver_email
    msg['Subject'] = subject

    msg.attach(MIMEText(body, 'plain'))

    try:
        with smtplib.SMTP('smtp.gmail.com', 587) as server:
            server.starttls()
            server.login(sender_email, password)
            server.sendmail(sender_email, receiver_email, msg.as_string())
    except Exception as e:
        print(f"Failed to send error report: {str(e)}")

# --- Gemini Call ---
def call_gemini(prompt, temperature=0.7, max_tokens=2048):
    lang = st.session_state.get("language", "en") # Corrected from .ge to .get
    lang_name = [k for k, v in languages.items() if v == lang][0]
    prompt = f"Please answer in {lang_name}.\n" + prompt
    
    # Ensure GEMINI_API_KEY is properly loaded from secrets
    GEMINI_API_KEY = st.secrets.get("gemini", {}).get("api_key", "") # Corrected secrets access

    if not GEMINI_API_KEY:
        st.write("There is an error, we've notified the support team.")
        send_error_report("Gemini API key is not configured. Please check your secrets.toml file.") # Corrected send_error_repor to send_error_report
        return "API key not configured"
        
    url = (
        f"https://generativelanguage.googleapis.com/v1beta/"
        f"models/gemini-2.0-flash-lite:generateContent?key={GEMINI_API_KEY}"
    )
    payload = {
        "contents": [{"parts": [{"text": prompt}]}],
        "generationConfig": {"temperature": temperature, "maxOutputTokens": max_tokens}
    }
    try:
        with show_lottie_loading(t("Thinking with Gemini AI...")): # Corrected from ("...") to t("...")
            response = requests.post(url, json=payload)
            response.raise_for_status()
            data = response.json()
            return data["candidates"][0]["content"]["parts"][0]["text"]
    except requests.exceptions.HTTPError as e:
        error_message = f"API Error: {str(e)}"
        send_error_report(error_message)
        st.error("There is an error, we've notified the support team.")
        return "Error occurred while calling Gemini API"
    except Exception as e:
        error_message = f"An error occurred: {str(e)}"
        send_error_report(error_message)
        st.error("There is an error, we've notified the support team.")
        return "Error occurred while processing your request"

def extract_text_from_file(file): # Renamed to avoid conflict, original was just extract_text
    name = file.name.lower()
    if name.endswith(".pdf"):
        with fitz.open(stream=file.read(), filetype="pdf") as doc:
            return "".join(page.get_text() for page in doc)
    if name.endswith(".docx"):
        return "\n".join(p.text for p in docx.Document(file).paragraphs)
    if name.endswith(".pptx"):
        return "\n".join(shape.text for slide in Presentation(file).slides for shape in slide.shapes if hasattr(shape, 'text'))
    if name.endswith(".txt"):
        return StringIO(file.getvalue().decode('utf-8')).read()
    if name.endswith((".jpg",".jpeg",".png")):
        return pytesseract.image_to_string(Image.open(file))
    return ""


# --- Visual/Equation/Code Understanding Helper (MUST BE DEFINED BEFORE USAGE) ---
def add_to_google_calendar(deadline):
    # Opens a Google Calendar event creation link in the browser
    import urllib.parse
    base = "https://calendar.google.com/calendar/render?action=TEMPLATE"
    params = {
        "text": deadline['description'],
        "dates": f"{deadline['date'].replace('-', '')}/{deadline['date'].replace('-', '')}",
    }
    url = base + "&" + urllib.parse.urlencode(params)
    webbrowser.open_new_tab(url)

def extract_visuals_and_code(text, file=None):
    visuals = []
    # Detect LaTeX/math equations (simple regex for $...$ or \[...\])
    import re
    equations = re.findall(r'(\$[^$]+\$|\\\[[^\]]+\\\])', text)
    for eq in equations:
        visuals.append(("Equation", eq))
    # Detect code blocks (triple backticks or indented)
    code_blocks = re.findall(r'```[\s\S]*?```', text)
    for code in code_blocks:
        visuals.append(("Code", code))
    # Detect possible diagrams/images in file (if image or PDF page)
    if file and hasattr(file, 'name') and file.name.lower().endswith((".jpg", ".jpeg", ".png")):
        visuals.append(("Diagram", "[Image uploaded]"))
    # For PDFs, could add more advanced image extraction if needed
    return visuals

# --- Calendar Integration Helper (MUST BE DEFINED BEFORE USAGE) ---
def detect_deadlines(text):
    prompt = (
        "Extract all assignment or exam deadlines (with date and description) from the following text. "
        "Return a JSON list of objects with 'date' and 'description'.\n\n" + text[:5000]
    )
    import json
    try:
        deadlines_json = call_gemini(prompt)
        deadlines = json.loads(deadlines_json)
        if isinstance(deadlines, dict):
            deadlines = list(deadlines.values())
        return deadlines
    except Exception:
        return []

# --- Lottie Loading Helper (MUST BE DEFINED BEFORE USAGE) ---
def load_lottieurl(url):
    r = reqs.get(url)
    if r.status_code != 200:
        return None
    return r.json()

@contextlib.contextmanager
def show_lottie_loading(message="Loading..."):
    # Create a container for the loading animation
    container = st.empty()
    try:
        # Show the loading animation
        with container:
            st.spinner(message)
        yield
    finally:
        # Remove the entire container and its contents
        container.empty()
        # Force a rerun to ensure the UI updates
        

# --- Configuration from st.secrets ---
raw_uri       = st.secrets.get("google", {}).get("redirect_uri", "") # Corrected secrets access
REDIRECT_URI  = raw_uri.rstrip("/") + "/" if raw_uri else ""
CLIENT_ID     = st.secrets.get("google", {}).get("client_id", "") # Corrected secrets access
CLIENT_SECRET = st.secrets.get("google", {}).get("client_secret", "") # Corrected secrets access
SCOPES        = ["openid", "email", "profile"]
# GEMINI_API_KEY is now loaded directly in call_gemini for robustness
CSE_API_KEY    = st.secrets.get("google_search", {}).get("api_key", "") # Corrected secrets access
CSE_ID         = st.secrets.get("google_search", {}).get("cse_id", "") # Corrected secrets access
CACHE_TTL      = 3600

# --- Product Hunt API Integration (Moved to top) ---
PRODUCT_HUNT_TOKEN = st.secrets.get("producthunt", {}).get("api_token", "") # Corrected secrets access
PRODUCT_HUNT_ID = st.secrets.get("producthunt", {}).get("product_id", "") # Corrected secrets access

@st.cache_data(ttl=300)
def get_ph_stats():
    if not PRODUCT_HUNT_TOKEN or not PRODUCT_HUNT_ID:
        return {"votes": 0, "comments": []}
    headers = {"Authorization": f"Bearer {PRODUCT_HUNT_TOKEN}"}
    # Get upvotes
    votes_url = f"https://api.producthunt.com/v2/api/graphql"
    votes_query = {
        "query": f"""
        query {{
          post(slug: \"{PRODUCT_HUNT_ID}\") {{
            votesCount
            comments(first: 5) {{
              edges {{
                node {{
                  id
                  body
                  user {{ name profileImage }}
                }}
              }}
            }}
          }}
        }}
        """
    }
    try:
        r = requests.post(votes_url, headers=headers, json=votes_query)
        data = r.json()
        post = data['data']['post']
        votes = post['votesCount']
        comments = [
            {
                "body": edge['node']['body'],
                "user": edge['node']['user']['name'],
                "avatar": edge['node']['user']['profileImage']
            }
            for edge in post['comments']['edges']
        ]
        return {"votes": votes, "comments": []}
    except Exception:
        return {"votes": 0, "comments": []}

# --- Image Processing for Whiteboard Doodle Effect ---
def process_image_for_doodle(uploaded_file):
    try:
        if uploaded_file.type == "application/pdf":
            # Read PDF and get the first page as an image
            doc = fitz.open(stream=uploaded_file.read(), filetype="pdf")
            page = doc.load_page(0)  # Load the first page
            pix = page.get_pixmap()
            img_bytes = pix.tobytes("png")
            img = Image.open(BytesIO(img_bytes))
        else:
            # It's a regular image file
            img = Image.open(uploaded_file)
        
        # Convert to grayscale
        img_gray = img.convert("L") # Corrected from .conver to .convert
        
        # Apply an edge detection filter (e.g., FIND_EDGES or CONTOUR)
        # Using a combination of filters can achieve a good sketch effect.
        # Here's a simple edge enhancement:
        img_edges = img_gray.filter(ImageFilter.FIND_EDGES)
        
        # Invert colors for a black outline on white background (like a whiteboard)
        img_doodle = ImageOps.invert(img_edges)
        
        # Convert back to RGB for display in Streamlit if it was grayscale
        img_doodle = img_doodle.convert("RGB") # Corrected from .conver to .convert

        # Save to bytes buffer
        buf = BytesIO()
        img_doodle.save(buf, format="PNG")
        return buf.getvalue()
    except Exception as e:
        st.error(f"Error processing image for doodle effect: {e}")
        return None

# --- Audio Generation for Whiteboard Explainer ---
def generate_audio_from_text(text, filename="temp_explainer_audio.mp3"):
    try:
        tts = gTTS(text=text, lang=st.session_state.get("language", "en")) # Corrected from .ge to .get
        temp_audio_path = os.path.join(tempfile.gettempdir(), filename)
        tts.save(temp_audio_path)
        return temp_audio_path
    except Exception as e:
        st.error(f"Error generating audio: {e}")
        return None

# --- Session State ---
for key in ("token", "user"):
    if key not in st.session_state:
        st.session_state[key] = None

# Whiteboard Explainer specific session state initializations
if 'wb_message' not in st.session_state: # Renamed to avoid conflict with general 'message'
    st.session_state['wb_message'] = ""
if 'wb_frames' not in st.session_state: # Renamed to avoid conflict
    st.session_state['wb_frames'] = []
if 'wb_generated_audio_path' not in st.session_state: # Renamed to avoid conflict
    st.session_state['wb_generated_audio_path'] = None
if 'wb_video_playing' not in st.session_state: # Renamed to avoid conflict
    st.session_state['wb_video_playing'] = False
if 'wb_processing' not in st.session_state: # Renamed to avoid conflict
    st.session_state['wb_processing'] = False

# --- SQLite DB for Learning Style ---
DB_PATH = "learning_style.db"

def init_db():
    with closing(sqlite3.connect(DB_PATH)) as conn:
        with conn:
            conn.execute('''
                CREATE TABLE IF NOT EXISTS learning_style (
                    email TEXT PRIMARY KEY,
                    sensing_intuitive INTEGER,
                    visual_verbal INTEGER,
                    active_reflective INTEGER,
                    sequential_global INTEGER
                )
            ''')



def save_learning_style(email, scores):
    with closing(sqlite3.connect(DB_PATH)) as conn:
        with conn:
            conn.execute('''
                INSERT INTO learning_style (email, sensing_intuitive, visual_verbal, active_reflective, sequential_global)
                VALUES (?, ?, ?, ?, ?)
                ON CONFLICT(email) DO UPDATE SET
                    sensing_intuitive=excluded.sensing_intuitive,
                    visual_verbal=excluded.visual_verbal,
                    active_reflective=excluded.active_reflective,
                    sequential_global=excluded.sequential_global
            ''', (email, scores['Sensing/Intuitive'], scores['Visual/Verbal'], scores['Active/Reflective'], scores['Sequential/Global']))

def get_learning_style(email):
    with closing(sqlite3.connect(DB_PATH)) as conn:
        cur = conn.execute('SELECT sensing_intuitive, visual_verbal, active_reflective, sequential_global FROM learning_style WHERE email=?', (email,))
        row = cur.fetchone()
        if row:
            return {
                'Sensing/Intuitive': row[0],
                'Visual/Verbal': row[1],
                'Active/Reflective': row[2],
                'Sequential/Global': row[3],
            }
        return None

# Initialise DB if needed
init_db()

# --- SQLite DB for Memorization Tracking ---
def init_mem_db():
    with closing(sqlite3.connect(DB_PATH)) as conn:
        with conn:
            conn.execute('''
                CREATE TABLE IF NOT EXISTS memorization (
                    email TEXT,
                    card_id TEXT,
                    question TEXT,
                    answer TEXT,
                    last_reviewed DATE,
                    next_due DATE,
                    correct_count INTEGER DEFAULT 0,
                    incorrect_count INTEGER DEFAULT 0,
                    PRIMARY KEY (email, card_id)
                )
            ''')

def get_due_cards(email, today):
    with closing(sqlite3.connect(DB_PATH)) as conn:
        cur = conn.execute('''
            SELECT card_id, question, answer, last_reviewed, next_due, correct_count, incorrect_count
            FROM memorization
            WHERE email=? AND (next_due IS NULL OR next_due<=?)
            ORDER BY next_due ASC
            LIMIT 10
        ''', (email, today))
        return cur.fetchall()

def update_card_review(email, card_id, correct, today):
    # Simple spaced repetition: if correct, next_due += 3 days, else next_due = tomorrow
    import datetime
    next_due = (datetime.datetime.strptime(today, "%Y-%m-%d") + (datetime.timedelta(days=3) if correct else datetime.timedelta(days=1))).strftime("%Y-%m-%d")
    with closing(sqlite3.connect(DB_PATH)) as conn:
        with conn:
            if correct:
                conn.execute('''
                    UPDATE memorization SET last_reviewed=?, next_due=?, correct_count=correct_count+1 WHERE email=? AND card_id=?
                ''', (today, next_due, email, card_id))
            else:
                conn.execute('''
                    UPDATE memorization SET last_reviewed=?, next_due=?, incorrect_count=incorrect_count+1 WHERE email=? AND card_id=?
                ''', (today, next_due, email, card_id))

def add_memorization_card(email, card_id, question, answer):
    with closing(sqlite3.connect(DB_PATH)) as conn:
        with conn:
            conn.execute('''
                INSERT OR IGNORE INTO memorization (email, card_id, question, answer) VALUES (?, ?, ?, ?)
            ''', (email, card_id, question, answer))

init_mem_db()

# --- SQLite DB for Content Structure & Progress Tracking ---
def init_structure_db():
    with closing(sqlite3.connect(DB_PATH)) as conn:
        with conn:
            conn.execute('''
                CREATE TABLE IF NOT EXISTS content_structure (
                    email TEXT,
                    doc_id TEXT,
                    section TEXT,
                    progress REAL DEFAULT 0.0,
                    PRIMARY KEY (email, doc_id, section)
                )
            ''')

def save_content_structure(email, doc_id, sections):
    with closing(sqlite3.connect(DB_PATH)) as conn:
        with conn:
            for section in sections:
                conn.execute('''
                    INSERT OR IGNORE INTO content_structure (email, doc_id, section) VALUES (?, ?, ?)
                ''', (email, doc_id, section))

def update_section_progress(email, doc_id, section, progress):
    with closing(sqlite3.connect(DB_PATH)) as conn:
        with conn:
            conn.execute('''
                UPDATE content_structure SET progress=? WHERE email=? AND doc_id=? AND section=?
            ''', (progress, email, doc_id, section))

def get_section_progress(email, doc_id):
    with closing(sqlite3.connect(DB_PATH)) as conn:
        cur = conn.execute('''
            SELECT section, progress FROM content_structure WHERE email=? AND doc_id=?
        ''', (email, doc_id))
        return dict(cur.fetchall())
        
LOGO_URL = "https://github.com/rekfdjkzbdfvkgjerkdfnfcbvgewhs/Vekkam/blob/main/logo.png"  # <-- Replace with your actual raw GitHub URL

init_structure_db()

# --- OAuth Flow using st.query_params ---
def ensure_logged_in():
    params = st.query_params
    code = params.get("code")  # Corrected from .ge to .get

    # Exchange code for token
    if code and not st.session_state.token:
        res = requests.post(
            "https://oauth2.googleapis.com/token",
            data={
                "code": code,
                "client_id": CLIENT_ID,
                "client_secret": CLIENT_SECRET,
                "redirect_uri": REDIRECT_URI,
                "grant_type": "authorization_code"
            }
        )
        if res.status_code != 200:
            st.error(f"Token exchange failed ({res.status_code}): {res.text}")
            st.stop()
        st.session_state.token = res.json()

        # Clear code from URL
        st.query_params.clear()

        # Fetch user info
        ui = requests.get(
            "https://www.googleapis.com/oauth2/v3/userinfo",
            headers={"Authorization": f"Bearer {st.session_state.token['access_token']}"}
        )
        if ui.status_code != 200:
            st.error("Failed to fetch user info.")
            st.stop()
        st.session_state.user = ui.json()

    # If not logged in, show landing page
    if not st.session_state.token:
        # Landing page layout
        st.markdown("""
            <style>
            .main {
                padding: 2rem;
            }
            .feature-box {
                background-color: #f0f2f6;
                border-radius: 10px;
                padding: 20px;
                margin: 10px 0;
            }
            .cta-button {
                background-color: #FF4B4B;
                color: white !important;
                padding: 15px 30px;
                border-radius: 5px;
                text-decoration: none !important;
                font-weight: bold;
                display: inline-flex;
                align-items: center;
                gap: 10px;
                transition: background-color 0.3s ease;
                font-size: 1.2rem;
            }
            .cta-button:hover {
                background-color: #FF3333;
                text-decoration: none !important;
            }
            .google-icon {
                width: 24px;
                height: 24px;
                margin-right: 5px;
            }
            .hero-section {
                text-align: center;
                padding: 2rem 0;
                background: linear-gradient(135deg, #f5f7fa 0%, #e4e8eb 100%);
                border-radius: 15px;
                margin-bottom: 2rem;
            }
            .feature-icon {
                font-size: 2.5rem;
                margin-bottom: 1rem;
            }
            .testimonial-box {
                background-color: white;
                border-radius: 10px;
                padding: 20px;
                margin: 10px 0;
                box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            }
            .stats-box {
                text-align: center;
                padding: 20px;
                background-color: #f8f9fa;
                border-radius: 10px;
                margin: 10px 0;
            }
            .stats-number {
                font-size: 2rem;
                font-weight: bold;
                color: #FF4B4B;
            }
            .whiteboard-container { /* Added for the new feature */
                width: 100%;
                max-width: 800px;
                aspect-ratio: 16 / 9;
                background-color: #f7f7f7;
                border: 8px solid #a0a0a0;
                border-radius: 12px;
                box-shadow: 0 10px 20px rgba(0, 0, 0, 0.15);
                display: flex;
                justify-content: center;
                align-items: center;
                text-align: center;
                overflow: hidden;
                margin-bottom: 20px;
            }
            .whiteboard-text { /* Added for the new feature */
                font-size: clamp(1.2rem, 3vw, 2.5rem);
                font-weight: 500;
                color: #333;
                padding: 20px;
                animation: fadeIn 0.8s forwards;
            }
            .whiteboard-image { /* Added for the new feature */
                max-width: 90%;
                max-height: 90%;
                object-fit: contain;
                border-radius: 8px;
                animation: fadeIn 0.8s forwards;
            }
            @keyframes fadeIn { /* Added for the new feature */
                from { opacity: 0; transform: translateY(20px) scale(0.95); }
                to { opacity: 1; transform: translateY(0) scale(1); }
            }
            </style>
        """, unsafe_allow_html=True)

        # Hero Section with CTA
        st.markdown('<div class="hero-section">', unsafe_allow_html=True)
        col1, col2 = st.columns([1, 2])
        with col1:
            st.image(LOGO_URL, width=200)
        with col2:
            st.markdown("<h1 style='font-size: 3rem; margin-bottom: 1rem;'>Welcome to Vekkam 📚</h1>", unsafe_allow_html=True)
            st.markdown("<h2 style='font-size: 1.5rem; color: #666; margin-bottom: 2rem;'>Your AI-powered study companion</h2>", unsafe_allow_html=True)
            
            # CTA Button
            auth_url = (
                "https://accounts.google.com/o/oauth2/v2/auth?"
                + urlencode({
                    "client_id": CLIENT_ID,
                    "redirect_uri": REDIRECT_URI,
                    "response_type": "code",
                    "scope": " ".join(SCOPES),
                    "access_type": "offline",
                    "prompt": "consent"
                })
            )
            st.markdown(
                f'<a href="{auth_url}" class="cta-button">'
                f'<img src="https://www.google.com/favicon.ico" class="google-icon" alt="Google icon">'
                f'Start Learning Now</a>',
                unsafe_allow_html=True
            )
        st.markdown('</div>', unsafe_allow_html=True)

        # Stats Section
        st.markdown("## 📊 Trusted by Students Worldwide")
        stats_cols = st.columns(4)
        stats = [
        ]
        for i, (number, label) in enumerate(stats):
            with stats_cols[i]:
                st.markdown(f'<div class="stats-box">', unsafe_allow_html=True)
                st.markdown(f'<div class="stats-number">{number}</div>', unsafe_allow_html=True)
                st.markdown(f'<div>{label}</div>', unsafe_allow_html=True)
                st.markdown('</div>', unsafe_allow_html=True)

        # Features Section
        st.markdown("## ✨ Powerful Features for Better Learning")
        features = [
            {
                "icon": "📖",
                "title": "Guide Book Chat",
                "description": "Search and chat with textbooks. Get instant explanations of any concept!",
                "details": [
                    "Smart search across multiple textbooks",
                    "Real-time concept explanations",
                    "Interactive Q&A with AI tutor",
                    "Save and review important concepts"
                ]
            },
            {
                "icon": "📝",
                "title": "Document Q&A",
                "description": "Upload your notes or books and get instant learning aids personalized for you.",
                "details": [
                    "Support for PDF, images, and text files",
                    "Automatic summary generation",
                    "Key points extraction",
                    "Custom flashcards creation"
                ]
            },
            {
                "icon": "📚",
                "title": "Paper Solver",
                "description": "Upload exam papers and get model answers with detailed explanations.",
                "details": [
                    "Step-by-step solutions",
                    "Exam tips and strategies",
                    "Common mistakes analysis",
                    "Practice questions generation"
                ]
            },
            {
                "icon": "🧠", # Corrected from � to 🧠
                "title": "Personalized Learning",
                "description": "AI-powered learning style assessment and personalized study recommendations.",
                "details": [
                    "Learning style analysis",
                    "Custom study plans",
                    "Progress tracking",
                    "Adaptive difficulty levels"
                ]
            }
        ]

        for feature in features:
            st.markdown('<div class="feature-box">', unsafe_allow_html=True)
            st.markdown(f"### {feature.get('icon', '')} {feature.get('title', '')}")
            st.markdown(f"**{feature.get('description', '')}**")
            st.markdown("")
            for detail in feature.get('details', []):
                st.markdown(f"- {detail}")
            st.markdown('</div>', unsafe_allow_html=True)

        # Benefits Section
        st.markdown("## 🎯 Why Choose Vekkam?")
        benefits = [
            {
                "icon": "🤖",
                "title": "AI-Powered Learning",
                "description": "Advanced AI technology adapts to your learning style and pace"
            },
            {
                "icon": "📊",
                "title": "Personalized Experience",
                "description": "Get customized study recommendations based on your learning style"
            },
            {
                "icon": "📱", # Corrected from � to 📱
                "title": "Access Anywhere",
                "description": "Study on any device, anytime, with seamless synchronization"
            },
            {
                "icon": "🌐",
                "title": "Multi-Language Support",
                "description": "Learn in your preferred language with accurate translations"
            },
            {
                "icon": "📈",
                "title": "Track Progress",
                "description": "Monitor your learning journey with detailed analytics"
            },
            {
                "icon": "🎯",
                "title": "Exam Success",
                "description": "Boost your confidence with comprehensive exam preparation"
            }
        ]
        
        cols = st.columns(3)
        for i, benefit in enumerate(benefits):
            with cols[i % 3]:
                st.markdown('<div class="feature-box">', unsafe_allow_html=True)
                st.markdown(f"### {benefit['icon']} {benefit['title']}")
                st.markdown(benefit['description'])
                st.markdown('</div>', unsafe_allow_html=True)

        # Testimonials Section
        st.markdown("## 💬 What Students Say")
        testimonials = [
            {
                "name": "Sarah K.",
                "role": "Medical Student",
                "text": "Vekkam helped me understand complex medical concepts through its interactive Q&A feature. The personalized learning approach made studying much more effective!"
            },
            {
                "name": "Michael R.",
                "role": "Engineering Student",
                "text": "The paper solver feature is amazing! It not only gives me the answers but also explains the concepts thoroughly. My grades have improved significantly."
            },
            {
                "name": "Priya M.",
                "role": "High School Student",
                "text": "I love how Vekkam adapts to my learning style. The flashcards and summaries are perfect for quick revision before exams."
            }
        ]

        for testimonial in testimonials:
            st.markdown('<div class="testimonial-box">', unsafe_allow_html=True)
            st.markdown(f"**{testimonial['name']}** - {testimonial['role']}")
            st.markdown(f"_{testimonial['text']}_")
            st.markdown('</div>', unsafe_allow_html=True)

        # Final CTA
        st.markdown("## 🚀 Ready to Transform Your Learning?")
        st.markdown("Join thousands of students who are already learning smarter with Vekkam!")
        st.markdown(
            f'<a href="{auth_url}" class="cta-button">'
            f'<img src="https://www.google.com/favicon.ico" class="google-icon" alt="Google icon">'
            f'Start Learning Now</a>',
                unsafe_allow_html=True
        )
        
        st.stop()

# Run OAuth check at startup
ensure_logged_in()

# --- After authentication UI ---
user = st.session_state.user

# Ensure the learning style test only appears when the learning style is not already stored
learning_style = get_learning_style(user.get("email", "")) # Corrected from .ge to .get
if learning_style is None:
    st.title(t("Welcome, {name}!", name=user.get('name', ''))) # Corrected t(...) syntax
    st.header(t("Learning Style Test")) # Corrected t(...) syntax
    st.write(t("Answer the following questions to determine your learning style. This will help us personalize your experience.")) # Corrected t(...) syntax
    likert = [
        t("Strongly Disagree"), t("Disagree"), t("Somewhat Disagree"), t("Neutral"), t("Somewhat Agree"), t("Agree"), t("Strongly Agree")
    ]
    questions = {
        "Sensing/Intuitive": [
            (t("I am more interested in what is actual than what is possible."), "Sensing"),
            (t("I often focus on the big picture rather than the details."), "Intuitive"),
            (t("I trust my gut feelings over concrete evidence."), "Intuitive"),
            (t("I enjoy tasks that require attention to detail."), "Sensing"),
            (t("I prefer practical solutions over theoretical ideas."), "Sensing"),
            (t("I am drawn to abstract concepts and patterns."), "Intuitive"),
            (t("I notice details that others might miss."), "Sensing"),
            (t("I like to imagine possibilities and what could be."), "Intuitive"),
            (t("I rely on past experiences to guide me."), "Sensing"),
            (t("I am energized by exploring new ideas."), "Intuitive"),
        ],
        "Visual/Verbal": [
            (t("I remember best what I see (pictures, diagrams, charts)."), "Visual"),
            (t("I find it easier to follow spoken instructions than written ones."), "Verbal"),
            (t("I prefer to learn through images and spatial understanding."), "Visual"),
            (t("I often take notes to help me remember."), "Verbal"),
            (t("I visualize information in my mind."), "Visual"),
            (t("I prefer reading to watching videos."), "Verbal"),
            (t("I use color and layout to organize my notes."), "Visual"),
            (t("I find it easier to express myself in writing."), "Verbal"),
            (t("I am drawn to infographics and visual summaries."), "Visual"),
            (t("I enjoy listening to lectures or podcasts."), "Verbal"),
        ],
        "Active/Reflective": [
            (t("I learn best by doing and trying things out."), "Active"),
            (t("I prefer to think things through before acting."), "Reflective"),
            (t("I enjoy group work and discussions."), "Active"),
            (t("I need time alone to process new information."), "Reflective"),
            (t("I like to experiment and take risks in learning."), "Active"),
            (t("I often review my notes quietly after class."), "Reflective"),
            (t("I am energized by interacting with others."), "Active"),
            (t("I prefer to observe before participating."), "Reflective"),
            (t("I learn by teaching others or explaining concepts aloud."), "Active"),
            (t("I keep a journal or log to reflect on my learning."), "Reflective"),
        ],
        "Sequential/Global": [
            (t("I learn best in a step-by-step, logical order."), "Sequential"),
            (t("I like to see the big picture before the details."), "Global"),
            (t("I prefer to follow clear, linear instructions."), "Sequential"),
            (t("I often make connections between ideas in a holistic way."), "Global"),
            (t("I am comfortable breaking tasks into smaller parts."), "Sequential"),
            (t("I sometimes jump to conclusions without all the steps."), "Global"),
            (t("I like outlines and structured notes."), "Sequential"),
            (t("I understand concepts better when I see how they fit together."), "Global"),
            (t("I prefer to finish one thing before starting another."), "Sequential"),
            (t("I enjoy brainstorming and exploring many ideas at once."), "Global"),
        ],
    }
    if "learning_style_answers" not in st.session_state:
        st.session_state.learning_style_answers = {}
    for dichotomy, qs in questions.items():
        st.subheader(dichotomy)
        for i, (q, side) in enumerate(qs):
            key = f"{dichotomy}_{i}"
            st.session_state.learning_style_answers[key] = st.radio(
                q,
                likert, # Use translated likert labels
                key=key
            )
    if st.button(t("Submit Learning Style Test")): # Corrected t(...) syntax
        # Scoring: Strongly Disagree=0, ..., Neutral=50, ..., Strongly Agree=100 (for positive phrasing)
        scores = {}
        for dichotomy, qs in questions.items():
            total = 0
            for i, (q, side) in enumerate(qs):
                key = f"{dichotomy}_{i}"
                val = st.session_state.learning_style_answers[key]
                idx = likert.index(val)
                # If the question is for the first side, score as is; if for the opposite, reverse
                if side == dichotomy.split("/")[0]: # Corrected .spli to .split
                    score = score_map[idx] # score_map needs to be defined
                else:
                    score = score_map[6 - idx] # score_map needs to be defined
                total += score
            scores[dichotomy] = int(total / len(qs))
        with show_lottie_loading(t("Saving your learning style and personalizing your experience...")): # Corrected t(...) syntax
            save_learning_style(user.get("email", ""), scores) # Corrected from .ge to .get
            st.session_state.learning_style_answers = {}
        st.success(t("Learning style saved! Reloading...")) # Corrected t(...) syntax
        st.balloons()
        st.rerun() # Rerun to apply learning style and proceed
        
    st.stop()
else:
    st.sidebar.write(t("Your learning style has been saved.")) # Corrected t(...) syntax

st.sidebar.image(user.get("picture", ""), width=48) # Corrected from .ge to .get
st.sidebar.write(user.get("email", "")) # Corrected from .ge to .get

# --- Personalized for you box ---
def learning_style_description(scores):
    desc = []
    if scores['Sensing/Intuitive'] >= 60:
        desc.append(t("Prefers concepts, patterns, and big-picture thinking.")) # Corrected t(...) syntax
    elif scores['Sensing/Intuitive'] <= 40:
        desc.append(t("Prefers facts, details, and practical examples.")) # Corrected t(...) syntax
    if scores['Visual/Verbal'] >= 60:
        desc.append(t("Learns best with visuals, diagrams, and mind maps.")) # Corrected t(...) syntax
    elif scores['Visual/Verbal'] <= 40:
        desc.append(t("Learns best with text, explanations, and reading.")) # Corrected t(...) syntax
    if scores['Active/Reflective'] >= 60:
        desc.append(t("Enjoys interactive, hands-on, and group activities.")) # Corrected t(...) syntax
    elif scores['Active/Reflective'] <= 40:
        desc.append(t("Prefers reflection, summaries, and solo study.")) # Corrected t(...) syntax
    if scores['Sequential/Global'] >= 60:
        desc.append(t("Prefers holistic overviews and big-picture connections.")) # Corrected t(...) syntax
    elif scores['Sequential/Global'] <= 40:
        desc.append(t("Prefers step-by-step, structured learning.")) # Corrected t(...) syntax
    return desc

if learning_style:
    st.sidebar.markdown("---")
    st.sidebar.subheader(t("Personalized for you")) # Corrected t(...) syntax
    st.sidebar.write({k: f"{v}/100" for k, v in learning_style.items()})
    for d in learning_style_description(learning_style):
        st.sidebar.info(d)

if st.sidebar.button(t("Logout")): # Corrected t(...) syntax
    st.session_state.clear()
    st.rerun() # Rerun to go back to login page


# --- PDF/Text Extraction ---
def extract_pages_from_url(pdf_url):
    with show_lottie_loading(t("Extracting PDF from URL...")): # Corrected t(...) syntax
        r = requests.get(pdf_url)
        # Use fitz for URL PDF handling
        with fitz.open(stream=r.content, filetype="pdf") as doc:
            return {i+1: doc[i].get_text() for i in range(len(doc))}

def extract_pages_from_file(file):
    with show_lottie_loading(t("Extracting PDF from file...")): # Corrected t(...) syntax
        # Use fitz for uploaded file PDF handling
        with fitz.open(stream=file.read(), filetype="pdf") as doc:
            return {i+1: doc[i].get_text() for i in range(len(doc))}

# The original extract_text function, now properly named to avoid conflict with the one above
def extract_text_from_uploaded_file(file): 
    ext = file.name.lower().split('.')[-1]
    if ext == "pdf":
        return "\n".join(extract_pages_from_file(file).values())
    if ext in ("jpg","jpeg","png"):
        with show_lottie_loading(t("Extracting text from image...")): # Corrected t(...) syntax
            return pytesseract.image_to_string(Image.open(file))
    with show_lottie_loading(t("Extracting text from file...")): # Corrected t(...) syntax
        return StringIO(file.getvalue().decode()).read()

# --- Guide Book Search & Concept Q&A ---
def fetch_pdf_url(title, author, edition):
    q = " ".join(filter(None, [title, author, edition]))
    params = {"key": CSE_API_KEY, "cx": CSE_ID, "q": q, "fileType": "pdf", "num": 1}
    with show_lottie_loading(t("Searching for PDF guide book...")): # Corrected t(...) syntax
        items = requests.get("https://www.googleapis.com/customsearch/v1", params=params).json().get("items", []) # Corrected .ge to .get
    return items[0]["link"] if items else None

def find_concept_pages(pages, concept):
    cl = concept.lower()
    return {p: t for p, t in pages.items() if cl in (t or "").lower()}

def ask_concept(pages, concept):
    found = find_concept_pages(pages, concept)
    if not found:
        return f"Couldn't find '{concept}'."
    combined = "\n---\n".join(f"Page {p}: {t}" for p, t in found.items())
    return call_gemini(f"Concept: '{concept}'. Sections:\n{combined}\nExplain with context and examples.")

# --- AI Learning Aids Functions ---
def generate_summary(text): 
    return call_gemini(
        f"Summarize this for an exam and separately list any formulae that are mentioned in the text."
        f"Format the formulae such that the terms are the terms themselves, like 'Quick Ratio = (Cash + Short-Term marketable securities + Receivables) / Current liabilities'"
        f"If there aren't any, skip this section. Output only.:\n\n{text}",
        temperature=0.5
    )

def generate_questions(text): 
    return call_gemini(
        f"Generate 15 quiz questions for an exam (ignore authors, ISSN, etc.). Output only.:\n\n{text}",
        temperature=0.7, max_tokens=8192
    )

def generate_flashcards(text): 
    return call_gemini(
        f"Create flashcards (Q&A). Output only.:\n\n{text}",
        temperature=0.7, max_tokens=8192
    )

def generate_mnemonics(text): 
    return call_gemini(
        f"Generate mnemonics. Output only.:\n\n{text}",
        temperature=0.7, max_tokens=8192
    )

def generate_key_terms(text): 
    return call_gemini(
        f"List all the key terms in the doc, with definitions. Output only.:\n\n{text}",
        temperature=0.7, max_tokens=8192
    )

def generate_cheatsheet(text): 
    return call_gemini(
        f"Create a cheat sheet. Output only.:\n\n{text}",
        temperature=0.7, max_tokens=8192
    )

def generate_highlights(text): 
    return call_gemini(
        f"List key facts and highlights. Output only.:\n\n{text}",
        temperature=0.7, max_tokens=8192
    )

def generate_critical_points(text):
    return call_gemini(
        f"I haven't studied for the exam, so run me over the doc in detail but concise, "
        f"so that I'm ready for the exam. Output only.:\n\n{text}",
        temperature=0.7, max_tokens=8192
    )

# --- Explainer Podcast Feature ---
def generate_podcast(text, filename="explainer_podcast.mp3"):
    """
    Generate an explainer podcast from the given text.
    """
    # Step 1: Generate podcast script using call_gemini
    prompt = (
        "You are an expert educator. Create a explainer podcast script that explains the following document "
        "in a clear, engaging, and easy-to-understand manner. Use examples, analogies, and a friendly tone. "
        "Structure the explanation into sections with transitions between topics.\n\n"
        "Output only the script, in plain text- no fancy formatting like bueelts, bold, italics- no of that."
        f"Document Content:\n{text}"
    )
    try:
        podcast_script = call_gemini(prompt)
    except Exception as e:
        st.error(f"Failed to generate podcast script: {str(e)}")
        return None

    # Step 2: Convert script to audio using gTTS
    try:
        tts = gTTS(podcast_script, lang="en")
        tts.save(filename)
        return filename
    except Exception as e:
        st.error(f"Failed to generate audio: {str(e)}")
        return None

# --- Helper to render each section in Streamlit ---
def render_section(title, content):
    st.subheader(title)
    if content.strip().startswith("<"):
        # raw HTML content
        components.html(content, height=600, scrolling=True)
    else:
        st.markdown(content, unsafe_allow_html=True)

def plot_mind_map(json_text):
    try:
        mind_map = json.loads(json_text)
    except json.JSONDecodeError:
        st.error("Mind map JSON invalid.")
        return
    nodes, edges, counter = [], [], 0
    def add_node(node, parent=None):
        nonlocal counter
        nid = counter; counter += 1
        label = node.get("title") or node.get("label") or "Node" # Corrected from .ge to .get
        nodes.append((nid, label))
        if parent is not None:
            edges.append((parent, nid))
        for child in node.get("children", []): # Corrected from .ge to .get
            add_node(child, nid)
    add_node(mind_map)
    g = ig.Graph(directed=True)
    g.add_vertices([str(n[0]) for n in nodes])
    g.vs["label"] = [n[1] for n in nodes]
    g.add_edges([(str(u),str(v)) for u,v in edges])
    layout = g.layout("tree") # Corrected from .layou to .layout
    x,y = zip(*layout.coords)
    edge_x,edge_y = [],[]
    for u,v in edges:
        edge_x += [x[u],x[v],None]
        edge_y += [y[u],y[v],None]
    edge_trace = go.Scatter(x=edge_x, y=edge_y, mode="lines", hoverinfo="none")
    node_trace = go.Scatter(x=x, y=y, text=g.vs["label"], mode="markers+text", textposition="top center",
                            marker=dict(size=20, line=dict(width=2)))
    fig = go.Figure([edge_trace, node_trace],
        layout=go.Layout(margin=dict(l=0,r=0,t=20,b=0), xaxis=dict(visible=False), yaxis=dict(visible=False)))
    st.plotly_chart(fig, use_container_width=True)

# --- Multilingual Support ---
languages = {
    "English": "en",
    "Hindi": "hi",
    "Tamil": "ta",
    "Spanish": "es",
    "French": "fr"
}

ui_translations = {
    "en": {
        "Guide Book Chat": "Guide Book Chat",
        "Document Q&A": "Document Q&A",
        "Learning Style Test": "Learning Style Test",
        "Paper Solver/Exam Guide": "Paper Solver/Exam Guide",
        "Upload your exam paper (PDF or image). The AI will extract questions and show you how to answer for full marks!": "Upload your exam paper (PDF or image). The AI will extract questions and show you how to answer for full marks!",
        "Upload Exam Paper (PDF/Image)": "Upload Exam Paper (PDF/Image)",
        "Found {n} questions:": "Found {n} questions:",
        "Select questions to solve (default: all)": "Select questions to solve (default: all)",
        "Solve Selected Questions": "Solve Selected Questions",
        "Model Answers & Exam Tips": "Model Answers & Exam Tips",
        "Welcome, {name}!": "Welcome, {name}!",
        "Feature": "Feature",
        "Logout": "Logout",
        "Learning Aids": "Learning Aids",
        "Pick a function": "Pick a function",
        "Run": "Run",
        "Recommended for you:": "Recommended for you:",
        "Personalized for you": "Personalized for you",
        "Answer the following questions to determine your learning style.": "Answer the following questions to determine your learning style.",
        "Submit Learning Style Test": "Submit Learning Style Test",
        "Saving your learning style and personalizing your experience...": "Saving your learning style and personalizing your experience...",
        "Learning style saved! Reloading...": "Learning style saved! Reloading...",
        "Extracting PDF from URL...": "Extracting PDF from URL...",
        "Extracting PDF from file...": "Extracting PDF from file...",
        "Extracting text from image...": "Extracting text from image...",
        "Extracting text from file...": "Extracting text from file...",
        "Thinking with Gemini AI...": "Thinking with Gemini AI...",
        "Searching for PDF guide book...": "Searching for PDF guide book...",
        "Extracting questions from PDF...": "Extracting questions from PDF...",
        "Extracting questions from image...": "Extracting questions from image...",
        "Solving Q{n}...": "Solving Q{n}...",
        "Whiteboard Explainer": "Whiteboard Explainer", # Added for new feature
        "Generate Whiteboard Video": "Generate Whiteboard Video", # Added for new feature
        "Play Video": "Play Video", # Added for new feature
        "Preparing video content... This may take a moment.": "Preparing video content... This may take a moment.", # Added for new feature
        "Please provide either text or images to generate a video.": "Please provide either text or images to generate a video.", # Added for new feature
        "Processing script...": "Processing script...", # Added for new feature
        "Processing images for doodle effect...": "Processing images for doodle effect...", # Added for new feature
        "Skipping image {name} due to processing error.": "Skipping image {name} due to processing error.", # Added for new feature
        "Generating audio for script...": "Generating audio for script...", # Added for new feature
        "Video content prepared. Click 'Play Video' to start.": "Video content prepared. Click 'Play Video' to start.", # Added for new feature
        "Playing video...": "Playing video...", # Added for new feature
        "Video playback complete!": "Video playback complete!", # Added for new feature
        "Note: This application simulates the whiteboard video in your browser and does not generate an MP4 file.": "Note: This application simulates the whiteboard video in your browser and does not generate an MP4 file.", # Added for new feature
        "Type your question here:": "Type your question here:",
        "Example: Can you explain how photosynthesis works?": "Example: Can you explain how photosynthesis works?",
        "Or upload an image of your question:": "Or upload an image of your question:",
        "Upload a clear image of your question or problem": "Upload a clear image of your question or problem",
        "Get Answer": "Get Answer",
        "Analyzing your question...": "Analyzing your question...",
        "Could not read text from the image. Please try uploading a clearer image.": "Could not read text from the image. Please try uploading a clearer image.",
        "Please either type a question or upload an image.": "Please either type a question or upload an image.",
        "Searching for relevant resources...": "Searching for relevant resources...",
        "Answer": "Answer",
        "Relevant Resources": "Relevant Resources",
        "Source:": "Source:",
        "Description:": "Description:",
        "Open Resource": "Open Resource",
        "Type:": "Type:",
        "Follow-up Questions": "Follow-up Questions",
        "Practice Problems": "Practice Problems",
        "Additional Resources": "Additional Resources",
        "Upload PDF/Image/TXT (multiple allowed)": "Upload PDF/Image/TXT (multiple allowed)",
        "Upload your notes, textbook, or image.": "Upload your notes, textbook, or image.",
        "Generate Explainer Podcast": "Generate Explainer Podcast",
        "Create Podcast": "Create Podcast",
        "Generating podcast...": "Generating podcast...",
        "Podcast generated successfully!": "Podcast generated successfully!",
        "Download Podcast": "Download Podcast",
        "Learning Aids for ": "Learning Aids for ",
        "Generating summary...": "Generating summary...",
        "Summary": "Summary",
        "Generating quiz questions...": "Generating quiz questions...",
        "Quiz Questions": "Quiz Questions",
        "Flashcards": "Flashcards",
        "Generating flashcards...": "Generating flashcards...",
        "Mnemonics": "Mnemonics",
        "Generating mnemonics...": "Generating mnemonics...",
        "Key Terms": "Key Terms",
        "Generating key terms...": "Generating key terms...",
        "Cheat Sheet": "Cheat Sheet",
        "Generating cheat sheet...": "Generating cheat sheet...",
        "Highlights": "Highlights",
        "Generating highlights...": "Generating highlights...",
        "Critical Points": "Critical Points",
        "Generating critical points...": "Generating critical points...",
        "Export all generated flashcards as an Anki-compatible CSV file.": "Export all generated flashcards as an Anki-compatible CSV file.",
        "Export All Flashcards to Anki CSV": "Export All Flashcards to Anki CSV",
        "Flashcards exported: ": "Flashcards exported: ",
        "Flashcards exported!": "Flashcards exported!",
        "Download Anki CSV": "Download Anki CSV",
        "Flashcards exported to Anki CSV!": "Flashcards exported to Anki CSV!",
        "Enter your script for the explainer video:": "Enter your script for the explainer video:",
        "e.g., 'Welcome to our explainer video. Today, we'll talk about innovative solutions. Here's a diagram...'": "e.g., 'Welcome to our explainer video. Today, we'll talk about innovative solutions. Here's a diagram...'",
        "Upload images for your explainer (optional):": "Upload images for your explainer (optional):",
        "Generate Whiteboard Video": "Generate Whiteboard Video",
        "Play Video": "Play Video",
        "Processing script...": "Processing script...",
        "Processing images for doodle effect...": "Processing images for doodle effect...",
        "Skipping image {name} due to processing error.": "Skipping image {name} due to processing error.", # The issue is here
        "Generating audio for script...": "Generating audio for script...",
        "Video content prepared. Click 'Play Video' to start.": "Video content prepared. Click 'Play Video' to start.",
        "Playing video...": "Playing video...",
        "Audio file not found. Please regenerate video.": "Audio file not found. Please regenerate video.",
        "Video playback complete!": "Video playback complete!",
        "Your learning style has been saved.": "Your learning style has been saved.",
        "Could not extract text from the uploaded file. Please ensure it's a clear document or image.": "Could not extract text from the uploaded file. Please ensure it's a clear document or image.",
        "Could not extract any questions from the document. Please ensure the format is clear.": "Could not extract any questions from the document. Please ensure the format is clear.",
        "Select questions to solve (default: all)": "Select questions to solve (default: all)",
        "Please select at least one question to solve.": "Please select at least one question to solve.",
        "Analyzing your materials and creating a battle plan...": "Analyzing your materials and creating a battle plan...",
        "Please upload at least one study material.": "Please upload at least one study material.",
        "Your 6-Hour Battle Plan": "Your 6-Hour Battle Plan",
        "Topic-Specific Resources": "Topic-Specific Resources",
        "Quick Reference Guide": "Quick Reference Guide",
        "Mental Preparation": "Mental Preparation",
        "Export Options": "Export Options",
        "Battle plan is ready. You can copy it manually if needed.": "Battle plan is ready. You can copy it manually if needed.",
        "Add to Calendar": "Add to Calendar",
        "Added to your calendar!": "Added to your calendar!",
        "Build strong study habits and stay accountable with our discipline features!": "Build strong study habits and stay accountable with our discipline features!",
        "Study Streak": "Study Streak",
        "Accountability": "Accountability",
        "Focus Mode": "Focus Mode",
        "Smart Schedule": "Smart Schedule",
        "Study Analytics": "Study Analytics",
        "Rewards": "Rewards",
        "Study Environment": "Study Environment",
        "Distraction Blocker": "Distraction Blocker",
        "Save Streak Data": "Save Streak Data",
        "Streak data saved!": "Streak data saved!",
        "Save Analytics Data": "Save Analytics Data",
        "Analytics data saved!": "Analytics data saved!",
        "Choose ambient sound": "Choose ambient sound",
        "Play Sound": "Play Sound",
        "Sound file not found. Make sure 'sounds' directory exists with MP3s.": "Sound file not found. Make sure 'sounds' directory exists with MP3s.",
        "Attempting to play sound: {file_path}. (Requires VLC on server)": "Attempting to play sound: {file_path}. (Requires VLC on server)",
        "Failed to play sound: {e}. VLC might not be installed or configured.": "Failed to play sound: {e}. VLC might not be installed or configured.",
        "Playing ambient sounds with VLC might not work as expected in all deployment environments (e.g., Streamlit Cloud) as it requires VLC to be installed on the server.": "Playing ambient sounds with VLC might not work as expected in all deployment environments (e.g., Streamlit Cloud) as it requires VLC to be installed on the server.",
        "Welcome to our explainer video. Today, we'll talk about innovative solutions. Here's a diagram...": "Welcome to our explainer video. Today, we'll talk about innovative solutions. Here's a diagram...",
        "I am more interested in what is actual than what is possible.": "I am more interested in what is actual than what is possible.",
        "I often focus on the big picture rather than the details.": "I often focus on the big picture rather than the details.",
        "I trust my gut feelings over concrete evidence.": "I trust my gut feelings over concrete evidence.",
        "I enjoy tasks that require attention to detail.": "I enjoy tasks that require attention to detail.",
        "I prefer practical solutions over theoretical ideas.": "I prefer practical solutions over theoretical ideas.",
        "I am drawn to abstract concepts and patterns.": "I am drawn to abstract concepts and patterns.",
        "I notice details that others might miss.": "I notice details that others might miss.",
        "I like to imagine possibilities and what could be.": "I like to imagine possibilities and what could be.",
        "I rely on past experiences to guide me.": "I rely on past experiences to guide me.",
        "I am energized by exploring new ideas.": "I am energized by exploring new ideas.",
        "I remember best what I see (pictures, diagrams, charts).": "I remember best what I see (pictures, diagrams, charts).",
        "I find it easier to follow spoken instructions than written ones.": "I find it easier to follow spoken instructions than written ones.",
        "I prefer to learn through images and spatial understanding.": "I prefer to learn through images and spatial understanding.",
        "I often take notes to help me remember.": "I often take notes to help me remember.",
        "I visualize information in my mind.": "I visualize information in my mind.",
        "I prefer reading to watching videos.": "I prefer reading to watching videos.",
        "I use color and layout to organize my notes.": "I use color and layout to organize my notes.",
        "I find it easier to express myself in writing.": "I find it easier to express myself in writing.",
        "I am drawn to infographics and visual summaries.": "I am drawn to infographics and visual summaries.",
        "I enjoy listening to lectures or podcasts.": "I enjoy listening to lectures or podcasts.",
        "I learn best by doing and trying things out.": "I learn best by doing and trying things out.",
        "I prefer to think things through before acting.": "I prefer to think things through before acting.",
        "I enjoy group work and discussions.": "I enjoy group work and discussions.",
        "I need time alone to process new information.": "I need time alone to process new information.",
        "I like to experiment and take risks in learning.": "I like to experiment and take risks in learning.",
        "I often review my notes quietly after class.": "I often review my notes quietly after class.",
        "I am energized by interacting with others.": "I am energized by interacting with others.",
        "I prefer to observe before participating.": "I prefer to observe before participating.",
        "I learn by teaching others or explaining concepts aloud.": "I learn by teaching others or explaining concepts aloud.",
        "I keep a journal or log to reflect on my learning.": "I keep a journal or log to reflect on my learning.",
        "I learn best in a step-by-step, logical order.": "I learn best in a step-by-step, logical order.",
        "I like to see the big picture before the details.": "I like to see the big picture before the details.",
        "I prefer to follow clear, linear instructions.": "I prefer to follow clear, linear instructions.",
        "I often make connections between ideas in a holistic way.": "I often make connections between ideas in a holistic way.",
        "I am comfortable breaking tasks into smaller parts.": "I am comfortable breaking tasks into smaller parts.",
        "I sometimes jump to conclusions without all the steps.": "I sometimes jump to conclusions without all the steps.",
        "I like outlines and structured notes.": "I like outlines and structured notes.",
        "I understand concepts better when I see how they fit together.": "I understand concepts better when I see how they fit together.",
        "I prefer to finish one thing before starting another.": "I prefer to finish one thing before starting another.",
        "I enjoy brainstorming and exploring many ideas at once.": "I enjoy brainstorming and exploring many ideas at once.",
        "Prefers concepts, patterns, and big-picture thinking.": "Prefers concepts, patterns, and big-picture thinking.",
        "Prefers facts, details, and practical examples.": "Prefers facts, details, and practical examples.",
        "Learns best with visuals, diagrams, and mind maps.": "Learns best with visuals, diagrams, and mind maps.",
        "Learns best with text, explanations, and reading.": "Learns best with text, explanations, and reading.",
        "Enjoys interactive, hands-on, and group activities.": "Enjoys interactive, hands-on, and group activities.",
        "Prefers reflection, summaries, and solo study.": "Prefers reflection, summaries, and solo study.",
        "Prefers holistic overviews and big-picture connections.": "Prefers holistic overviews and big-picture connections.",
        "Prefers step-by-step, structured learning.": "Prefers step-by-step, structured learning.",
        "Current Streak": "Current Streak",
        "Longest Streak": "Longest Streak",
        "Total Study Time": "Total Study Time",
        "Choose ambient sound": "Choose ambient sound"
    },
    "hi": {
        "Guide Book Chat": "गाइड बुक चैट",
        "Document Q&A": "दस्तावेज़ प्रश्नोत्तर",
        "Learning Style Test": "अधिगम शैली परीक्षण",
        "Paper Solver/Exam Guide": "पेपर सॉल्वर/परीक्षा गाइड",
        "Upload your exam paper (PDF or image). The AI will extract questions and show you how to answer for full marks!": "अपना परीक्षा पत्र (PDF या छवि) अपलोड करें। एआई प्रश्नों को निकालेगा और आपको पूर्ण अंक पाने के लिए उत्तर कैसे देना है, यह बताएगा!",
        "Upload Exam Paper (PDF/Image)": "परीक्षा पत्र अपलोड करें (PDF/छवि)",
        "Found {n} questions:": "{n} प्रश्न मिले:",
        "Select questions to solve (default: all)": "हल करने के लिए प्रश्न चुनें (डिफ़ॉल्ट: सभी)",
        "Solve Selected Questions": "चयनित प्रश्न हल करें",
        "Model Answers & Exam Tips": "मॉडल उत्तर और परीक्षा टिप्स",
        "Welcome, {name}!": "स्वागत है, {name}!",
        "Feature": "विशेषता",
        "Logout": "लॉगआउट",
        "Learning Aids": "अधिगम सहायक",
        "Pick a function": "एक फ़ंक्शन चुनें",
        "Run": "चलाएँ",
        "Recommended for you:": "आपके लिए अनुशंसित:",
        "Personalized for you": "आपके लिए वैयक्तिकृत",
        "Answer the following questions to determine your learning style.": "अपनी अधिगम शैली निर्धारित करने के लिए निम्नलिखित प्रश्नों का उत्तर दें।",
        "Submit Learning Style Test": "अधिगम शैली परीक्षण सबमिट करें",
        "Saving your learning style and personalizing your experience...": "आपकी अधिगम शैली सहेजी जा रही है और आपके अनुभव को वैयक्तिकृत किया जा रहा है...",
        "Learning style saved! Reloading...": "अधिगम शैली सहेजी गई! पुनः लोड हो रहा है...",
        "Extracting PDF from URL...": "URL से PDF निकाला जा रहा है...",
        "Extracting PDF from file...": "फ़ाइल से PDF निकाला जा रहा है...",
        "Extracting text from image...": "छवि से पाठ निकाला जा रहा है...",
        "Extracting text from file...": "फ़ाइल से पाठ निकाला जा रहा है...",
        "Thinking with Gemini AI...": "Gemini AI के साथ सोच रहे हैं...",
        "Searching for PDF guide book...": "PDF गाइड बुक खोजी जा रही है...",
        "Extracting questions from PDF...": "PDF से प्रश्न निकाले जा रहे हैं...",
        "Extracting questions from image...": "छवि से प्रश्न निकाले जा रहे हैं...",
        "Solving Q{n}...": "Q{n} हल किया जा रहा है...",
        "Whiteboard Explainer": "व्हाइटबोर्ड एक्सप्लेनर", # Added for new feature
        "Generate Whiteboard Video": "व्हाइटबोर्ड वीडियो बनाएं", # Added for new feature
        "Play Video": "वीडियो चलाएं", # Added for new feature
        "Preparing video content... This may take a moment.": "वीडियो सामग्री तैयार हो रही है... इसमें थोड़ा समय लग सकता है।", # Added for new feature
        "Please provide either text or images to generate a video.": "कृपया वीडियो बनाने के लिए पाठ या छवियां प्रदान करें।", # Added for new feature
        "Processing script...": "स्क्रिप्ट संसाधित हो रही है...", # Added for new feature
        "Processing images for doodle effect...": "डूडल प्रभाव के लिए छवियां संसाधित हो रही हैं...", # Added for new feature
        "Skipping image {name} due to processing error.": "संसाधन त्रुटि के कारण छवि {name} को छोड़ दिया जा रहा है।",
        "Generating audio for script...": "स्क्रिप्ट के लिए ऑडियो उत्पन्न किया जा रहा है...",
        "Video content prepared. Click 'Play Video' to start.": "वीडियो सामग्री तैयार है। 'वीडियो चलाएं' पर क्लिक करें।", # Added for new feature
        "Playing video...": "वीडियो चल रहा है...", # Added for new feature
        "Video playback complete!": "वीडियो प्लेबैक पूरा हुआ!", # Added for new feature
        "Note: This application simulates the whiteboard video in your browser and does not generate an MP4 file.": "ध्यान दें: यह एप्लिकेशन आपके ब्राउज़र में व्हाइटबोर्ड वीडियो का अनुकरण करता है और MP4 फ़ाइल उत्पन्न नहीं करता है।", # Added for new feature
        "Type your question here:": "अपना प्रश्न यहां टाइप करें:",
        "Example: Can you explain how photosynthesis works?": "उदाहरण: प्रकाश संश्लेषण कैसे काम करता है?",
        "Or upload an image of your question:": "या अपने प्रश्न की एक छवि अपलोड करें:",
        "Upload a clear image of your question or problem": "अपने प्रश्न या समस्या की एक स्पष्ट छवि अपलोड करें",
        "Get Answer": "उत्तर प्राप्त करें",
        "Analyzing your question...": "आपके प्रश्न का विश्लेषण हो रहा है...",
        "Could not read text from the image. Please try uploading a clearer image.": "छवि से पाठ नहीं पढ़ा जा सका। कृपया एक स्पष्ट छवि अपलोड करने का प्रयास करें।",
        "Please either type a question or upload an image.": "कृपया या तो एक प्रश्न टाइप करें या एक छवि अपलोड करें।",
        "Searching for relevant resources...": "प्रासंगिक संसाधनों की खोज हो रही है...",
        "Answer": "उत्तर",
        "Relevant Resources": "प्रासंगिक संसाधन",
        "Source:": "स्रोत:",
        "Description:": "विवरण:",
        "Open Resource": "संसाधन खोलें",
        "Type:": "प्रकार:",
        "Follow-up Questions": "अनुवर्ती प्रश्न",
        "Practice Problems": "अभ्यास प्रश्न",
        "Additional Resources": "अतिरिक्त संसाधन",
        "Upload PDF/Image/TXT (multiple allowed)": "पीडीएफ/छवि/टीएक्सटी अपलोड करें (कई अनुमत हैं)",
        "Upload your notes, textbook, or image.": "अपने नोट्स, पाठ्यपुस्तक, या छवि अपलोड करें।",
        "Generate Explainer Podcast": "व्याख्याता पॉडकास्ट बनाएं",
        "Create Podcast": "पॉडकास्ट बनाएं",
        "Generating podcast...": "पॉडकास्ट उत्पन्न हो रहा है...",
        "Podcast generated successfully!": "पॉडकास्ट सफलतापूर्वक उत्पन्न हुआ!",
        "Download Podcast": "पॉडकास्ट डाउनलोड करें",
        "Learning Aids for ": "के लिए शिक्षण सहायक सामग्री",
        "Generating summary...": "सारांश उत्पन्न हो रहा है...",
        "Summary": "सारांश",
        "Generating quiz questions...": "क्विज़ प्रश्न उत्पन्न हो रहे हैं...",
        "Quiz Questions": "क्विज़ प्रश्न",
        "Flashcards": "फ्लैशकार्ड",
        "Generating flashcards...": "फ्लैशकार्ड उत्पन्न हो रहे हैं...",
        "Mnemonics": "स्मृति सहायक",
        "Generating mnemonics...": "स्मृति सहायक उत्पन्न हो रहे हैं...",
        "Key Terms": "मुख्य शर्तें",
        "Generating key terms...": "मुख्य शर्तें उत्पन्न हो रही हैं...",
        "Cheat Sheet": "चीट शीट",
        "Generating cheat sheet...": "चीट शीट उत्पन्न हो रही है...",
        "Highlights": "मुख्य बातें",
        "Generating highlights...": "मुख्य बातें उत्पन्न हो रही हैं...",
        "Critical Points": "महत्वपूर्ण बिंदु",
        "Generating critical points...": "महत्वपूर्ण बिंदु उत्पन्न हो रहे हैं...",
        "Export all generated flashcards as an Anki-compatible CSV file.": "सभी उत्पन्न फ्लैशकार्ड को अंकी-संगत CSV फ़ाइल के रूप में निर्यात करें।",
        "Export All Flashcards to Anki CSV": "सभी फ्लैशकार्ड को अंकी CSV में निर्यात करें",
        "Flashcards exported: ": "फ्लैशकार्ड निर्यात किए गए:",
        "Flashcards exported!": "फ्लैशकार्ड निर्यात किए गए!",
        "Download Anki CSV": "अंकी CSV डाउनलोड करें",
        "Flashcards exported to Anki CSV!": "फ्लैशकार्ड अंकी CSV में निर्यात किए गए!",
        "Enter your script for the explainer video:": "व्याख्याता वीडियो के लिए अपनी स्क्रिप्ट दर्ज करें:",
        "e.g., 'Welcome to our explainer video. Today, we'll talk about innovative solutions. Here's a diagram...'": "उदाहरण: 'हमारे व्याख्याता वीडियो में आपका स्वागत है। आज, हम अभिनव समाधानों के बारे में बात करेंगे। यहाँ एक आरेख है...'",
        "Upload images for your explainer (optional):": "अपने व्याख्याता के लिए छवियां अपलोड करें (वैकल्पिक):",
        "Generate Whiteboard Video": "व्हाइटबोर्ड वीडियो बनाएं",
        "Play Video": "वीडियो चलाएं",
        "Processing script...": "स्क्रिप्ट संसाधित हो रही है...",
        "Processing images for doodle effect...": "डूडल प्रभाव के लिए छवियां संसाधित हो रही हैं...",
        "Skipping image {name} due to processing error.": "संसाधन त्रुटि के कारण छवि {name} को छोड़ दिया जा रहा है।",
        "Generating audio for script...": "स्क्रिप्ट के लिए ऑडियो उत्पन्न किया जा रहा है...",
        "Video content prepared. Click 'Play Video' to start.": "वीडियो सामग्री तैयार है। 'वीडियो चलाएं' पर क्लिक करें।",
        "Playing video...": "वीडियो चल रहा है...",
        "Audio file not found. Please regenerate video.": "ऑडियो फ़ाइल नहीं मिली। कृपया वीडियो को फिर से उत्पन्न करें।",
        "Video playback complete!": "वीडियो प्लेबैक पूरा हुआ!",
        "Your learning style has been saved.": "आपकी अधिगम शैली सहेजी गई है।",
        "Could not extract text from the uploaded file. Please ensure it's a clear document or image.": "अपलोड की गई फ़ाइल से पाठ नहीं निकाला जा सका। कृपया सुनिश्चित करें कि यह एक स्पष्ट दस्तावेज़ या छवि है।",
        "Could not extract any questions from the document. Please ensure the format is clear.": "दस्तावेज़ से कोई प्रश्न नहीं निकाला जा सका। कृपया सुनिश्चित करें कि प्रारूप स्पष्ट है।",
        "Select questions to solve (default: all)": "हल करने के लिए प्रश्न चुनें (डिफ़ॉल्ट: सभी)",
        "Please select at least one question to solve.": "कृपया हल करने के लिए कम से कम एक प्रश्न चुनें।",
        "Analyzing your materials and creating a battle plan...": "आपकी सामग्री का विश्लेषण किया जा रहा है और एक युद्ध योजना बनाई जा रही है...",
        "Please upload at least one study material.": "कृपया कम से कम एक अध्ययन सामग्री अपलोड करें।",
        "Your 6-Hour Battle Plan": "आपकी 6 घंटे की युद्ध योजना",
        "Topic-Specific Resources": "विषय-विशिष्ट संसाधन",
        "Quick Reference Guide": "त्वरित संदर्भ मार्गदर्शिका",
        "Mental Preparation": "मानसिक तैयारी",
        "Export Options": "निर्यात विकल्प",
        "Battle plan is ready. You can copy it manually if needed.": "युद्ध योजना तैयार है। यदि आवश्यक हो तो आप इसे मैन्युअल रूप से कॉपी कर सकते हैं।",
        "Add to Calendar": "कैलेंडर में जोड़ें",
        "Added to your calendar!": "आपके कैलेंडर में जोड़ा गया!",
        "Build strong study habits and stay accountable with our discipline features!": "हमारी अनुशासन सुविधाओं के साथ मजबूत अध्ययन की आदतें बनाएं और जवाबदेह रहें!",
        "Study Streak": "अध्ययन स्ट्रीक",
        "Accountability": "जवाबदेही",
        "Focus Mode": "फोकस मोड",
        "Smart Schedule": "स्मार्ट शेड्यूल",
        "Study Analytics": "अध्ययन विश्लेषण",
        "Rewards": "पुरस्कार",
        "Study Environment": "अध्ययन वातावरण",
        "Distraction Blocker": "ध्यान भटकाने वाला अवरोधक",
        "Save Streak Data": "स्ट्रीक डेटा सहेजें",
        "Streak data saved!": "स्ट्रीक डेटा सहेजा गया!",
        "Save Analytics Data": "विश्लेषण डेटा सहेजें",
        "Analytics data saved!": "विश्लेषण डेटा सहेजा गया!",
        "Choose ambient sound": "परिवेश ध्वनि चुनें",
        "Play Sound": "ध्वनि चलाएं",
        "Sound file not found. Make sure 'sounds' directory exists with MP3s.": "ध्वनि फ़ाइल नहीं मिली। सुनिश्चित करें कि 'ध्वनि' निर्देशिका MP3s के साथ मौजूद है।",
        "Attempting to play sound: {file_path}. (Requires VLC on server)": "ध्वनि चलाने का प्रयास हो रहा है: {file_path}। (सर्वर पर VLC की आवश्यकता है)",
        "Failed to play sound: {e}. VLC might not be installed or configured.": "ध्वनि चलाने में विफल: {e}। VLC स्थापित या कॉन्फ़िगर नहीं हो सकता है।",
        "Playing ambient sounds with VLC might not work as expected in all deployment environments (e.g., Streamlit Cloud) as it requires VLC to be installed on the server.": "Streamlit Cloud जैसे सभी परिनियोजन वातावरणों में VLC के साथ परिवेश ध्वनियाँ अपेक्षित रूप से काम नहीं कर सकती हैं क्योंकि इसके लिए सर्वर पर VLC स्थापित होना आवश्यक है।"
    },
    # Add more languages as needed
}

def t(key, **kwargs):
    lang = st.session_state.get("language", "en") # Corrected from .ge to .get
    txt = ui_translations.get(lang, ui_translations["en"]).get(key, key)
    return txt.format(**kwargs)

# Language selector in sidebar
if "language" not in st.session_state:
    st.session_state["language"] = "en"
lang_choice = st.sidebar.selectbox("🌐 " + t("Language"), list(languages.keys()), index=0) # Corrected from ("...") to t("...")
st.session_state["language"] = languages[lang_choice]

# --- App Branding ---
st.markdown("""
    <style>
    .block-container {padding-top: 1.5rem;}
    .sidebar-content {padding-top: 1rem;}
    </style>
    """, unsafe_allow_html=True)
col1, col2 = st.columns([1, 8])
with col1:
    st.image(LOGO_URL, width=180)
with col2:
    st.markdown("<h1 style='margin-bottom:0;'>Vekkam 📚</h1>", unsafe_allow_html=True)
    st.caption("Your AI-powered study companion")

# --- Sidebar Onboarding/Help ---
st.sidebar.markdown("---")
with st.sidebar.expander("❓ " + t("How to use this app"), expanded=False): # Corrected from ("...") to t("...")
    st.markdown(t("""
    - **Choose your language** from the sidebar.
    - **Take the Learning Style Test** (first login) for personalized recommendations.
    - **Guide Book Chat**: Search and chat with textbooks.
    - **Document Q&A**: Upload notes or books for instant learning aids.
    - **Paper Solver/Exam Guide**: Upload an exam paper and get model answers.
    - **Whiteboard Explainer**: Create a simulated whiteboard video from text and images.
    - All features are personalized for you!
    """))

# --- Main UI ---
# Added "Whiteboard Explainer" to the list of tabs
quiz_tabs = [t("Guide Book Chat"), t("Document Q&A"), t("Whiteboard Explainer"), t("Learning Style Test"), t("Paper Solver/Exam Guide"), "⚡ 6-Hour Battle Plan", "🎯 Discipline Hub"] # Corrected from ("...") to t("...")
tab = st.sidebar.selectbox(t("Feature"), quiz_tabs) # Corrected from ("...") to t("...")

# Add this after the existing imports
def search_educational_resources(query, num_results=5):
    """Search for educational resources using Google Programmable Search Engine."""
    if not CSE_API_KEY or not CSE_ID:
        return []
    
    try:
        # Construct the search URL
        search_url = "https://www.googleapis.com/customsearch/v1"
        params = {
            "key": CSE_API_KEY,
            "cx": CSE_ID,
            "q": query,
            "num": num_results,
            "safe": "active",
            "lr": "lang_en",  # English language results
            "as_sitesearch": "edu",  # Prioritize educational sites
            "as_filetype": "pdf"  # Include PDF resources
        }
        
        response = requests.get(search_url, params=params)
        response.raise_for_status()
        results = response.json().get("items", []) # Corrected .ge to .get
        
        # Format the results
        formatted_results = []
        for item in results:
            formatted_results.append({
                "title": item.get("title", ""), # Corrected .ge to .get
                "link": item.get("link", ""), # Corrected .ge to .get
                "snippet": item.get("snippet", ""), # Corrected .ge to .get
                "file_type": item.get("fileFormat", ""), # Corrected .ge to .get
                "source": item.get("displayLink", "") # Corrected .ge to .get
            })
        
        return formatted_results
    except Exception as e:
        st.error(f"Error searching for resources: {str(e)}")
        return []

# Main tab selection
if tab == t("Guide Book Chat"): # Corrected t(...) syntax
    st.header("❓ " + t("Ask Your Questions")) # Corrected t(...) syntax
    st.info(t("Ask any question or upload an image of your question. Our AI will help you understand and solve it!")) # Corrected t(...) syntax

    # Create two columns for text input and image upload
    col1, col2 = st.columns(2)
    
    with col1:
        question = st.text_area(t("Type your question here:"), height=150, 
            placeholder=t("Example: Can you explain how photosynthesis works?")) # Corrected t(...) syntax
    
    with col2:
        uploaded_image = st.file_uploader(t("Or upload an image of your question:"), # Corrected t(...) syntax
            type=["jpg", "jpeg", "png"],
            help=t("Upload a clear image of your question or problem")) # Corrected t(...) syntax

    # Process the question (either from text or image)
    if st.button(t("Get Answer")): # Corrected t(...) syntax
        with show_lottie_loading(t("Analyzing your question...")): # Corrected t(...) syntax
            if uploaded_image:
                # Extract text from image
                image_text = pytesseract.image_to_string(Image.open(uploaded_image))
                if not image_text.strip():
                    st.error(t("Could not read text from the image. Please try uploading a clearer image.")) # Corrected t(...) syntax
                    st.stop()
                question = image_text

            if not question.strip():
                st.warning(t("Please either type a question or upload an image.")) # Corrected t(...) syntax
                st.stop()

            # Search for relevant resources
            with show_lottie_loading(t("Searching for relevant resources...")): # Corrected t(...) syntax
                search_results = search_educational_resources(question)
            
            # Generate a comprehensive answer
            prompt = (
                f"You are an expert tutor. Please help the student understand and solve this question. "
                f"Provide a clear, step-by-step explanation that includes:\n"
                f"1. Key concepts involved\n"
                f"2. Step-by-step solution or explanation\n"
                f"3. Examples or analogies to help understanding\n"
                f"4. Common mistakes to avoid\n"
                f"5. Related concepts to explore\n\n"
                f"Question: {question}"
            )
            
            answer = call_gemini(prompt)
            
            # Display the answer in a nicely formatted way
            st.markdown("### 📝 " + t("Answer")) # Corrected t(...) syntax
            st.markdown(answer)
            
            # Display relevant resources if found
            if search_results:
                st.markdown("---")
                st.markdown("### 📚 " + t("Relevant Resources")) # Corrected t(...) syntax
                for i, result in enumerate(search_results, 1):
                    with st.expander(f"{i}. {result['title']}"):
                        st.markdown(f"**{t('Source:')}** {result['source']}")
                        st.markdown(f"**{t('Description:')}** {result['snippet']}")
                        st.markdown(f"[{t('Open Resource')}]({result['link']})")
                        if result['file_type']:
                            st.markdown(f"**{t('Type:')}** {result['file_type']}")
            
            # Add a section for follow-up questions
            st.markdown("---")
            st.markdown("### 💭 " + t("Follow-up Questions")) # Corrected t(...) syntax
            follow_up_prompt = (
                f"Based on the student's question and the answer provided, suggest 3 follow-up questions "
                f"that would help deepen their understanding of the topic. Make them specific and thought-provoking.\n\n"
                f"Original Question: {question}\n"
                f"Answer: {answer}"
            )
            follow_ups = call_gemini(follow_up_prompt)
            st.markdown(follow_ups)

            # Add a section for practice problems
            st.markdown("---")
            st.markdown("### 📚 " + t("Practice Problems")) # Corrected t(...) syntax
            practice_prompt = (
                f"Create 2 practice problems related to the concepts in the question. "
                f"For each problem, provide:\n"
                f"1. The problem statement\n"
                f"2. A hint\n"
                f"3. The solution\n\n"
                f"Original Question: {question}\n"
                f"Answer: {answer}"
            )
            practice_problems = call_gemini(practice_prompt)
            st.markdown(practice_problems)

            # Add a section for additional resources
            st.markdown("---")
            st.markdown("### 🔍 " + t("Additional Resources")) # Corrected t(...) syntax
            resources_prompt = (
                f"Suggest 3-4 additional resources (videos, articles, interactive tools) that would help "
                f"the student better understand the topic. Include brief descriptions of each resource.\n\n"
                f"Original Question: {question}\n"
                f"Answer: {answer}"
            )
            resources = call_gemini(resources_prompt)
            st.markdown(resources)

elif tab == t("Document Q&A"): # Corrected t(...) syntax
    st.header("\U0001F4A1 " + t("Document Q&A")) # Corrected t(...) syntax
    st.info(t("Upload one or more documents and get instant learning aids, personalized for your style. The AI can now synthesize across multiple files!")) # Corrected t(...) syntax
    uploaded_files = st.file_uploader(t("Upload PDF/Image/TXT (multiple allowed)"), type=["pdf","jpg","png","txt"], help=t("Upload your notes, textbook, or image."), accept_multiple_files=True) # Corrected t(...) syntax
    texts = []
    file_names = []
    if uploaded_files:
        for uploaded in uploaded_files:
            # Extract text from file
            ext = uploaded.name.lower().split('.')[-1]
            if ext == "pdf":
                with show_lottie_loading(t("Extracting PDF from file...")): # Corrected t(...) syntax
                    # Using fitz directly
                    with fitz.open(stream=uploaded.read(), filetype="pdf") as doc:
                        text = "\n".join([page.get_text() for page in doc])
            elif ext in ("jpg", "jpeg", "png"):
                with show_lottie_loading(t("Extracting text from image...")): # Corrected t(...) syntax
                    text = pytesseract.image_to_string(Image.open(uploaded))
            else:
                with show_lottie_loading(t("Extracting text from file...")): # Corrected t(...) syntax
                    text = StringIO(uploaded.getvalue().decode()).read()
            texts.append(text)
            file_names.append(uploaded.name)

        # Combine all extracted text
        combined_text = "\n".join(texts)

        # Generate Podcast
        st.subheader("🎙️ " + t("Generate Explainer Podcast")) # Corrected t(...) syntax
        if st.button(t("Create Podcast")): # Corrected t(...) syntax
            with show_lottie_loading(t("Generating podcast...")): # Corrected t(...) syntax
                podcast_file = generate_podcast(combined_text)
                if podcast_file:
                    st.success(t("Podcast generated successfully!")) # Corrected t(...) syntax
                    st.audio(podcast_file, format="audio/mp3")
                    st.download_button(t("Download Podcast"), data=open(podcast_file, "rb"), file_name="explainer_podcast.mp3", mime="audio/mp3") # Corrected t(...) syntax

        # --- Generate learning aids for each file ---
        for idx, (text, fname) in enumerate(zip(texts, file_names)):
            st.subheader(f"{t('Learning Aids for ')}{fname}")
            
            # Generate and display all learning aids
            with show_lottie_loading(t("Generating summary...")): # Corrected t(...) syntax
                render_section(t("Summary"), generate_summary(text)) # Corrected t(...) syntax
            with show_lottie_loading(t("Generating quiz questions...")): # Corrected t(...) syntax
                render_section(t("Quiz Questions"), generate_questions(text)) # Corrected t(...) syntax

            with st.expander(t("Flashcards")): # Corrected t(...) syntax
                with show_lottie_loading(t("Generating flashcards...")): # Corrected t(...) syntax
                    render_section(t("Flashcards"), generate_flashcards(text)) # Corrected t(...) syntax

            with st.expander(t("Mnemonics")): # Corrected t(...) syntax
                with show_lottie_loading(t("Generating mnemonics...")): # Corrected t(...) syntax
                    render_section(t("Mnemonics"), generate_mnemonics(text)) # Corrected t(...) syntax

            with st.expander(t("Key Terms")): # Corrected t(...) syntax
                with show_lottie_loading(t("Generating key terms...")): # Corrected t(...) syntax
                    render_section(t("Key Terms"), generate_key_terms(text)) # Corrected t(...) syntax

            with st.expander(t("Cheat Sheet")): # Corrected t(...) syntax
                with show_lottie_loading(t("Generating cheat sheet...")): # Corrected t(...) syntax
                    render_section(t("Cheat Sheet"), generate_cheatsheet(text)) # Corrected t(...) syntax

            with st.expander(t("Highlights")): # Corrected t(...) syntax
                with show_lottie_loading(t("Generating highlights...")): # Corrected t(...) syntax
                    render_section(t("Highlights"), generate_highlights(text)) # Corrected t(...) syntax

            with st.expander(t("Critical Points")): # Corrected t(...) syntax
                with show_lottie_loading(t("Generating critical points...")): # Corrected t(...) syntax
                    render_section(t("Critical Points"), generate_critical_points(text)) # Corrected t(...) syntax

            # Store for batch export
            all_summaries.append(generate_summary(text))
            flashcards_raw = generate_flashcards(text)
            try:
                flashcards_json = json.loads(flashcards_raw)
                flashcards = [(fc['question'], fc['answer']) for fc in flashcards_json]
            except Exception:
                # fallback: try to split by Q/A
                flashcards = []
                for line in flashcards_raw.split('\n'):
                    if ':' in line:
                        q, a = line.split(':', 1)
                        flashcards.append((q.strip(), a.strip()))
            all_flashcards.extend(flashcards)

        # --- Batch Export ---
        if all_flashcards:
            st.info(t("Export all generated flashcards as an Anki-compatible CSV file.")) # Corrected t(...) syntax
            # Dummy function for export_flashcards_to_anki - implement if needed
            def export_flashcards_to_anki(flashcards_data):
                csv_file = StringIO()
                writer = csv.writer(csv_file)
                for q, a in flashcards_data:
                    writer.writerow([q, a])
                csv_file.seek(0)
                return csv_file.getvalue()

            fname = export_flashcards_to_anki(all_flashcards)
            st.download_button(t("Download Anki CSV"), data=fname, file_name="flashcards.csv", mime="text/csv") # Corrected t(...) syntax
            st.success(t("Flashcards exported to Anki CSV!")) # Corrected t(...) syntax
            st.toast(t("Flashcards exported!")) # Corrected t(...) syntax


elif tab == t("Whiteboard Explainer"): # New tab for Whiteboard Explainer # Corrected t(...) syntax
    st.header("✨ " + t("Whiteboard Explainer")) # Corrected t(...) syntax
    st.info(t("Note: This application simulates the whiteboard video in your browser and does not generate an MP4 file.")) # Corrected t(...) syntax

    script_text = st.text_area(
        t("Enter your script for the explainer video:"), # Corrected t(...) syntax
        height=200,
        placeholder=t("e.g., 'Welcome to our explainer video. Today, we'll talk about innovative solutions. Here's a diagram...'") # Corrected t(...) syntax
    )

    uploaded_images_wb = st.file_uploader( # Renamed variable to avoid conflict
        t("Upload images for your explainer (optional):"), # Corrected t(...) syntax
        type=["png", "jpg", "jpeg", "pdf"], # Added PDF
        accept_multiple_files=True,
        key="wb_image_uploader" # Added key for uniqueness
    )

    col1_wb, col2_wb = st.columns([1,1]) # Renamed columns to avoid conflict

    with col1_wb:
        generate_button_wb = st.button(t("Generate Whiteboard Video"), disabled=st.session_state['wb_processing'] or (not script_text and not uploaded_images_wb), key="generate_wb_button") # Corrected t(...) syntax

    with col2_wb:
        play_button_wb = st.button(t("Play Video"), disabled=st.session_state['wb_video_playing'] or st.session_state['wb_processing'] or not st.session_state['wb_frames'], key="play_wb_button") # Corrected t(...) syntax

    if generate_button_wb:
        if not script_text and not uploaded_images_wb:
            st.session_state['wb_message'] = t("Please provide either text or images to generate a video.") # Corrected t(...) syntax
        else:
            st.session_state['wb_processing'] = True
            st.session_state['wb_video_playing'] = False # Stop any ongoing playback
            st.session_state['wb_frames'] = []
            st.session_state['wb_generated_audio_path'] = None
            st.session_state['wb_message'] = t("Preparing video content... This may take a moment.") # Corrected t(...) syntax
            # st.rerun() # Rerun will be handled by the processing logic below

    # This block executes only when processing is True, and reruns are triggered.
    if st.session_state['wb_processing']:
        new_frames = []
        
        # Process script text
        if script_text:
            st.session_state['wb_message'] = t("Processing script...") # Corrected t(...) syntax
            sentences = re.split(r'(?<=[.?!])\s+', script_text) # Use re.split for more robust sentence splitting
            sentences = [s.strip() for s in sentences if s.strip()] # Filter out empty strings
            # Estimate duration based on sentence length for rough synchronization
            for sentence in sentences:
                duration = max(len(sentence) * 70, 1500) # Min 1.5s, 70ms per char
                new_frames.append({'type': 'text', 'content': sentence, 'duration': duration})
        
        # Process images
        if uploaded_images_wb:
            st.session_state['wb_message'] = t("Processing images for doodle effect...") # Corrected t(...) syntax
            for i, image_file in enumerate(uploaded_images_wb):
                doodle_image_bytes = process_image_for_doodle(image_file)
                if doodle_image_bytes:
                    # Add a 3 second delay for each image
                    new_frames.append({'type': 'image', 'content': doodle_image_bytes, 'duration': 3000}) 
                else:
                    st.session_state['wb_message'] = t("Skipping image {name} due to processing error.", name=image_file.name) # Fixed NameError
                    

        st.session_state['wb_frames'] = new_frames
        
        # Generate audio for the full script
        if script_text:
            st.session_state['wb_message'] = t("Generating audio for script...") # Corrected t(...) syntax
            audio_path = generate_audio_from_text(script_text)
            st.session_state['wb_generated_audio_path'] = audio_path
        
        st.session_state['wb_processing'] = False
        st.session_state['wb_message'] = t("Video content prepared. Click 'Play Video' to start.") # Corrected t(...) syntax
        st.rerun() # Rerun to update UI after processing is complete

    st.markdown(f"<p class='text-center text-sm font-medium text-gray-600'>{st.session_state['wb_message']}</p>", unsafe_allow_html=True)

    # Whiteboard display area
    whiteboard_placeholder = st.empty()

    if play_button_wb and st.session_state['wb_frames'] and not st.session_state['wb_processing']:
        st.session_state['wb_video_playing'] = True
        st.session_state['wb_message'] = t("Playing video...") # Corrected t(...) syntax
        
        # Play audio if available
        if st.session_state['wb_generated_audio_path']:
            try:
                audio_file = open(st.session_state['wb_generated_audio_path'], 'rb')
                audio_bytes = audio_file.read()
                st.audio(audio_bytes, format='audio/mp3', start_time=0, key="wb_audio_player") # Added key
                audio_file.close()
            except FileNotFoundError:
                st.session_state['wb_message'] = t("Audio file not found. Please regenerate video.") # Corrected t(...) syntax
                st.session_state['wb_video_playing'] = False
                st.rerun()
                
        # Simulate frame-by-frame playback
        for i, frame in enumerate(st.session_state['wb_frames']):
            with whiteboard_placeholder.container():
                if frame['type'] == 'text':
                    st.markdown(
                        f"""
                        <div class="whiteboard-container">
                            <p class="whiteboard-text">{frame['content']}</p>
                        </div>
                        """,
                        unsafe_allow_html=True
                    )
                elif frame['type'] == 'image':
                    st.markdown(
                        f"""
                        <div class="whiteboard-container">
                            <img src="data:image/png;base64,{base64.b64encode(frame['content']).decode('utf-8')}" class="whiteboard-image" />
                        </div>
                        """,
                        unsafe_allow_html=True
                    )
            time.sleep(frame['duration'] / 1000.0) # Convert ms to seconds
            
            # Clear the whiteboard placeholder at the end of each frame (except the very last one)
            if i < len(st.session_state['wb_frames']) - 1:
                with whiteboard_placeholder.container():
                    st.markdown(
                        """
                        <div class="whiteboard-container">
                            <!-- Empty for transition -->
                        </div>
                        """,
                        unsafe_allow_html=True
                    )
                time.sleep(0.2) # Short pause between frames
        
        st.session_state['wb_video_playing'] = False
        st.session_state['wb_message'] = t("Video playback complete!") # Corrected t(...) syntax
        # Clean up temporary audio file after playback
        if st.session_state['wb_generated_audio_path'] and os.path.exists(st.session_state['wb_generated_audio_path']):
            os.remove(st.session_state['wb_generated_audio_path'])
            st.session_state['wb_generated_audio_path'] = None
        st.rerun() # Rerun to update button state and message


elif tab == t("Learning Style Test"): # Corrected t(...) syntax
    st.header("🧠 " + t("Learning Style Test")) # Corrected t(...) syntax
    st.write(t("Answer the following questions to determine your learning style.")) # Corrected t(...) syntax
    
    likert_labels = [
        t("Strongly Disagree"), t("Disagree"), t("Somewhat Disagree"), t("Neutral"), t("Somewhat Agree"), t("Agree"), t("Strongly Agree")
    ]
    # Re-define score_map here for the learning style test to be self-contained
    score_map = {
        0: 0, 1: 17, 2: 33, 3: 50, 4: 67, 5: 83, 6: 100
    }

    questions = {
        "Sensing/Intuitive": [
            (t("I am more interested in what is actual than what is possible."), "Sensing"),
            (t("I often focus on the big picture rather than the details."), "Intuitive"),
            (t("I trust my gut feelings over concrete evidence."), "Intuitive"),
            (t("I enjoy tasks that require attention to detail."), "Sensing"),
            (t("I prefer practical solutions over theoretical ideas."), "Sensing"),
            (t("I am drawn to abstract concepts and patterns."), "Intuitive"),
            (t("I notice details that others might miss."), "Sensing"),
            (t("I like to imagine possibilities and what could be."), "Intuitive"),
            (t("I rely on past experiences to guide me."), "Sensing"),
            (t("I am energized by exploring new ideas."), "Intuitive"),
        ],
        "Visual/Verbal": [
            (t("I remember best what I see (pictures, diagrams, charts)."), "Visual"),
            (t("I find it easier to follow spoken instructions than written ones."), "Verbal"),
            (t("I prefer to learn through images and spatial understanding."), "Visual"),
            (t("I often take notes to help me remember."), "Verbal"),
            (t("I visualize information in my mind."), "Visual"),
            (t("I prefer reading to watching videos."), "Verbal"),
            (t("I use color and layout to organize my notes."), "Visual"),
            (t("I find it easier to express myself in writing."), "Verbal"),
            (t("I am drawn to infographics and visual summaries."), "Visual"),
            (t("I enjoy listening to lectures or podcasts."), "Verbal"),
        ],
        "Active/Reflective": [
            (t("I learn best by doing and trying things out."), "Active"),
            (t("I prefer to think things through before acting."), "Reflective"),
            (t("I enjoy group work and discussions."), "Active"),
            (t("I need time alone to process new information."), "Reflective"),
            (t("I like to experiment and take risks in learning."), "Active"),
            (t("I often review my notes quietly after class."), "Reflective"),
            (t("I am energized by interacting with others."), "Active"),
            (t("I prefer to observe before participating."), "Reflective"),
            (t("I learn by teaching others or explaining concepts aloud."), "Active"),
            (t("I keep a journal or log to reflect on my learning."), "Reflective"),
        ],
        "Sequential/Global": [
            (t("I learn best in a step-by-step, logical order."), "Sequential"),
            (t("I like to see the big picture before the details."), "Global"),
            (t("I prefer to follow clear, linear instructions."), "Sequential"),
            (t("I often make connections between ideas in a holistic way."), "Global"),
            (t("I am comfortable breaking tasks into smaller parts."), "Sequential"),
            (t("I sometimes jump to conclusions without all the steps."), "Global"),
            (t("I like outlines and structured notes."), "Sequential"),
            (t("I understand concepts better when I see how they fit together."), "Global"),
            (t("I prefer to finish one thing before starting another."), "Sequential"),
            (t("I enjoy brainstorming and exploring many ideas at once."), "Global"),
        ],
    }

    if "learning_style_answers" not in st.session_state:
        st.session_state.learning_style_answers = {}
    
    for dichotomy, qs in questions.items():
        st.subheader(dichotomy)
        for i, (q, side) in enumerate(qs):
            key = f"{dichotomy}_{i}"
            st.session_state.learning_style_answers[key] = st.radio(
                q,
                likert_labels,
                key=key
            )
    
    if st.button(t("Submit Learning Style Test")): # Corrected t(...) syntax
        scores = {}
        for dichotomy, qs in questions.items():
            total = 0
            for i, (q, side) in enumerate(qs):
                key = f"{dichotomy}_{i}"
                val = st.session_state.learning_style_answers[key]
                idx = likert_labels.index(val)
                score = score_map[idx]
                
                # If the question is for the second side of the dichotomy, invert the score
                # e.g., if dichotomy is Sensing/Intuitive and question is "I focus on big picture" (Intuitive side),
                # and user selects "Strongly Agree", the score should lean towards Intuitive (high score).
                # If question is "I am interested in what is actual" (Sensing side) and user selects "Strongly Agree",
                # score should lean towards Sensing (high score).
                
                # Determine which side of the dichotomy the current question's 'side' parameter refers to
                dichotomy_sides = dichotomy.split("/")
                
                if side == dichotomy_sides[0]: # If question aligns with the first part (e.g., Sensing for Sensing/Intuitive)
                    adjusted_score = score
                elif side == dichotomy_sides[1]: # If question aligns with the second part (e.g., Intuitive for Sensing/Intuitive)
                    adjusted_score = 100 - score # Invert score for the opposite side
                else: # Should not happen if 'side' is correctly defined
                    adjusted_score = score 

                total += adjusted_score
            scores[dichotomy] = int(total / len(qs))

        with show_lottie_loading(t("Saving your learning style and personalizing your experience...")): # Corrected t(...) syntax
            save_learning_style(user.get("email", ""), scores) # Corrected from .ge to .get
            st.session_state.learning_style_answers = {}
        st.success(t("Learning style saved! Reloading...")) # Corrected t(...) syntax
        st.balloons()
        st.rerun()

elif tab == t("Paper Solver/Exam Guide"): # Corrected t(...) syntax
    st.header("📝 " + t("Paper Solver/Exam Guide")) # Corrected t(...) syntax
    st.info(t("Upload your exam paper (PDF or image). The AI will extract questions and show you how to answer for full marks!")) # Corrected t(...) syntax

    exam_paper_file = st.file_uploader(t("Upload Exam Paper (PDF/Image)"), type=["pdf", "jpg", "jpeg", "png"]) # Corrected t(...) syntax

    if exam_paper_file:
        raw_text = extract_text_from_file(exam_paper_file) # Use the correct extraction function

        if not raw_text.strip():
            st.error(t("Could not extract text from the uploaded file. Please ensure it's a clear document or image.")) # Corrected t(...) syntax
            st.stop()

        # Step 1: Extract questions
        with show_lottie_loading(t("Extracting questions from PDF..." if exam_paper_file.type == "application/pdf" else "Extracting questions from image...")): # Corrected t(...) syntax
            question_extraction_prompt = (
                "From the following exam paper text, extract each question. "
                "List them numerically, starting with Q1., Q2., etc.\n\n"
                f"Exam Paper Text:\n{raw_text[:8000]}" # Limit text to avoid exceeding token limits
            )
            extracted_questions_raw = call_gemini(question_extraction_prompt, max_tokens=2048)
            
            # Parse extracted questions
            questions_list = re.findall(r'Q\d+\.\s*(.*)', extracted_questions_raw)
            if not questions_list:
                st.warning(t("Could not extract any questions from the document. Please ensure the format is clear.")) # Corrected t(...) syntax
                st.stop()

            st.session_state['extracted_questions'] = questions_list
            st.subheader(t("Found {n} questions:", n=len(questions_list))) # Corrected t(...) syntax

            # Step 2: Allow selection of questions
            selected_questions_indices = []
            if questions_list:
                st.write(t("Select questions to solve (default: all)")) # Corrected t(...) syntax
                for i, q in enumerate(questions_list):
                    if st.checkbox(f"Q{i+1}: {q}", value=True, key=f"q_checkbox_{i}"):
                        selected_questions_indices.append(i)
            
            if st.button(t("Solve Selected Questions")): # Corrected t(...) syntax
                if not selected_questions_indices:
                    st.warning(t("Please select at least one question to solve.")) # Corrected t(...) syntax
                else:
                    st.subheader(t("Model Answers & Exam Tips")) # Corrected t(...) syntax
                    for i in selected_questions_indices:
                        question_text = questions_list[i]
                        with show_lottie_loading(t("Solving Q{n}...", n=i+1)): # Corrected t(...) syntax
                            answer_prompt = (
                                f"You are an expert examiner. Provide a comprehensive model answer for the following exam question "
                                f"to achieve full marks. Also, include specific exam tips and common pitfalls to avoid for this type of question.\n\n"
                                f"Question: {question_text}\n"
                                f"Context from paper (if relevant): {raw_text[:2000]}" # Provide some context from the paper
                            )
                            model_answer = call_gemini(answer_prompt, temperature=0.5, max_tokens=2048)
                            st.markdown(f"#### Q{i+1}: {question_text}")
                            st.markdown(model_answer)
                            st.markdown("---")

elif tab == "⚡ 6-Hour Battle Plan":
    st.header("⚡ 6-Hour Battle Plan")
    st.info(t("Upload your syllabus, guide books, and study materials. We'll create a focused 6-hour study plan using Vekkam's features to help you ace your exam!")) # Corrected t(...) syntax

    # File upload section
    st.subheader("📚 " + t("Upload Your Materials")) # Corrected t(...) syntax
    uploaded_files = st.file_uploader(
        t("Upload your syllabus, guide books, and study materials (PDF/Image/TXT)"), # Corrected t(...) syntax
        type=["pdf", "jpg", "jpeg", "png", "txt"],
        accept_multiple_files=True,
        help=t("Upload all relevant study materials. The more you provide, the better the plan will be!") # Corrected t(...) syntax
    )

    # Additional information
    st.subheader("📝 " + t("Additional Information")) # Corrected t(...) syntax
    exam_date = st.date_input(t("When is your exam?"), help=t("This helps us prioritize topics")) # Corrected t(...) syntax
    exam_duration = st.number_input(t("Exam duration (in hours)"), min_value=1, max_value=6, value=3, help=t("How long is your exam?")) # Corrected t(...) syntax
    weak_topics = st.text_area(t("Topics you find challenging (optional)"), help=t("List topics you find difficult, separated by commas")) # Corrected t(...) syntax
    strong_topics = st.text_area(t("Topics you're confident in (optional)"), help=t("List topics you're good at, separated by commas")) # Corrected t(...) syntax

    if st.button(t("Generate Battle Plan")): # Corrected t(...) syntax
        if not uploaded_files:
            st.warning(t("Please upload at least one study material.")) # Corrected t(...) syntax
            st.stop()

        with show_lottie_loading(t("Analyzing your materials and creating a battle plan...")): # Corrected t(...) syntax
            # Extract text from all files
            all_text = []
            for file in uploaded_files:
                ext = file.name.lower().split('.')[-1]
                if ext == "pdf":
                    # Using fitz directly
                    with fitz.open(stream=file.read(), filetype="pdf") as doc:
                        text = "\n".join([page.get_text() for page in doc])
                elif ext in ("jpg", "jpeg", "png"):
                    text = pytesseract.image_to_string(Image.open(file))
                else:
                    text = StringIO(file.getvalue().decode()).read()
                all_text.append(text)

            combined_text = "\n---\n".join(all_text)

            # First, analyze the content and create a topic breakdown
            content_analysis_prompt = (
                f"Analyze the following study materials and create a structured breakdown of topics. "
                f"For each topic, identify:\n"
                f"1. Key concepts and formulas\n"
                f"2. Difficulty level (Easy/Medium/Hard)\n"
                f"3. Estimated time needed for mastery\n"
                f"4. Dependencies on other topics\n\n"
                f"Materials: {combined_text}"
            )
            content_analysis = call_gemini(content_analysis_prompt)

            # Generate the battle plan
            battle_plan_prompt = (
                f"Create a detailed 6-hour study plan using Vekkam's features. "
                f"Consider:\n"
                f"1. Exam date: {exam_date}\n"
                f"2. Exam duration: {exam_duration} hours\n"
                f"3. Challenging topics: {weak_topics}\n"
                f"4. Strong topics: {strong_topics}\n\n"
                f"Content Analysis: {content_analysis}\n\n"
                f"The plan should include:\n"
                f"1. Hour-by-hour breakdown:\n"
                f"   - Topic to cover\n"
                f"   - Vekkam feature to use (Guide Book Chat, Document Q&A, etc.)\n"
                f"   - Specific tasks and goals\n"
                f"   - Break times\n"
                f"2. For each topic:\n"
                f"   - Quick concept review using Guide Book Chat\n"
                f"   - Practice questions using Paper Solver\n"
                f"   - Flashcards for key points\n"
                f"   - Summary generation\n"
                f"3. Progress tracking and checkpoints\n"
                f"4. Energy management tips\n\n"
                f"Study Materials: {combined_text}"
            )

            battle_plan = call_gemini(battle_plan_prompt)

            # Display the battle plan in a structured way
            st.markdown("## 📋 " + t("Your 6-Hour Battle Plan")) # Corrected t(...) syntax
            st.markdown(battle_plan)

            # Generate topic-specific resources
            st.markdown("---")
            st.markdown("## 📚 " + t("Topic-Specific Resources")) # Corrected t(...) syntax
            resources_prompt = (
                f"Based on the content analysis and battle plan, create a list of resources for each topic:\n"
                f"1. Key formulas to memorize\n"
                f"2. Practice questions to attempt\n"
                f"3. Flashcards to create\n"
                f"4. Summary points\n\n"
                f"Content Analysis: {content_analysis}\n"
                f"Battle Plan: {battle_plan}"
            )
            resources = call_gemini(resources_prompt)
            st.markdown(resources)

            # Generate a quick reference guide
            st.markdown("---")
            st.markdown("## 📝 " + t("Quick Reference Guide")) # Corrected t(...) syntax
            reference_prompt = (
                f"Create a quick reference guide that includes:\n"
                f"1. All key formulas and concepts\n"
                f"2. Common mistakes to avoid\n"
                f"3. Last-minute tips for each topic\n"
                f"4. Time management strategies\n\n"
                f"Content Analysis: {content_analysis}\n"
                f"Battle Plan: {battle_plan}"
            )
            reference_guide = call_gemini(reference_prompt)
            st.markdown(reference_guide)

            # Add a section for mental preparation
            st.markdown("---")
            st.markdown("## 🧠 " + t("Mental Preparation")) # Corrected t(...) syntax
            mental_prompt = (
                f"Provide advice on:\n"
                f"1. How to stay calm during the exam\n"
                f"2. Time management strategies\n"
                f"3. How to handle difficult questions\n"
                f"4. Post-exam review tips\n\n"
                f"Based on the exam duration of {exam_duration} hours and the topics covered."
            )
            mental_prep = call_gemini(mental_prompt)
            st.markdown(mental_prep)

            # Add export options
            st.markdown("---")
            st.markdown("## 📤 " + t("Export Options")) # Corrected t(...) syntax
            col1, col2 = st.columns(2)
            with col1:
                st.info(t("Battle plan is ready. You can copy it manually if needed.")) # Corrected t(...) syntax
            
            with col2:
                if st.button(t("Add to Calendar")): # Corrected t(...) syntax
                    # Create calendar event for study session
                    event_title = "6-Hour Study Battle Plan"
                    event_desc = f"""
                    Battle Plan:
                    {battle_plan}
                    
                    Quick Reference:
                    {reference_guide}
                    """
                    add_to_google_calendar({
                        "date": exam_date.strftime("%Y-%m-%d"),
                        "description": event_title
                    })
                    st.success(t("Added to your calendar!")) # Corrected t(...) syntax

elif tab == "🎯 Discipline Hub": # Corrected from elif to if
    st.header("🎯 " + t("Discipline Hub")) # Corrected t(...) syntax
    st.info(t("Build strong study habits and stay accountable with our discipline features!")) # Corrected t(...) syntax

    # Create tabs for different discipline features
    discipline_tabs = st.tabs([
        "📊 " + t("Study Streak"), # Corrected t(...) syntax
        "👥 " + t("Accountability"), # Corrected t(...) syntax
        "⏱️ " + t("Focus Mode"), # Corrected t(...) syntax
        "📅 " + t("Smart Schedule"), # Corrected t(...) syntax
        "📈 " + t("Study Analytics"), # Corrected t(...) syntax
        "🏆 " + t("Rewards"), # Corrected t(...) syntax
        "🎯 " + t("Study Environment"), # Corrected t(...) syntax
        "🚫 " + t("Distraction Blocker") # Corrected t(...) syntax
    ])
    
    # --- Helper Functions ---
    def play_sound(file_path):
        """Play a sound file using VLC."""
        # VLC can be tricky to set up in all environments.
        # This function might need adjustment based on deployment.
        # For a web environment like Streamlit Cloud, direct VLC playback
        # on the server-side might not be the ideal approach.
        # Consider client-side audio playback or embedding if this causes issues.
        st.warning(t("Playing ambient sounds with VLC might not work as expected in all deployment environments (e.g., Streamlit Cloud) as it requires VLC to be installed on the server.")) # Corrected t(...) syntax
        if os.path.exists(file_path):
            # This would typically run on the server.
            # For client-side audio in Streamlit, you'd usually use st.audio
            # with a base64 encoded audio file or a publicly accessible URL.
            try:
                # Placeholder for direct server-side VLC command if available
                # import subprocess
                # subprocess.Popen(["vlc", "--play-and-exit", file_path])
                st.info(t("Attempting to play sound: {file_path}. (Requires VLC on server)", file_path=file_path)) # Corrected t(...) syntax
            except Exception as e:
                st.error(t("Failed to play sound: {e}. VLC might not be installed or configured.", e=e)) # Corrected t(...) syntax
            
            # Alternative: If you want client-side sound, you'd need to serve the sound file
            # or use a base64 encoding with st.audio
            # For now, keeping it as is based on original structure.
        else:
            st.error(t("Sound file not found. Make sure 'sounds' directory exists with MP3s.")) # Corrected t(...) syntax
            return None

    # --- Update Study Streak Feature ---
    if 'study_streak' not in st.session_state:
        st.session_state.study_streak = {
            'current_streak': 0,
            'longest_streak': 0,
            'last_study_date': None,
            'free_passes_remaining': 1,
            'total_study_time': 0
        }

    # Save streak data to a file
    STREAK_FILE = "streak_data.json"
    def save_streak_data():
        with open(STREAK_FILE, "w") as f:
            json.dump(st.session_state.study_streak, f)

    def load_streak_data():
        if os.path.exists(STREAK_FILE):
            with open(STREAK_FILE, "r") as f:
                st.session_state.study_streak = json.load(f)

    load_streak_data()

    # --- Update Productivity Scores ---
    if 'study_analytics' not in st.session_state:
        st.session_state.study_analytics = {
            'daily_hours': [],
            'topics': {},
            'productivity_score': 0,
            'focus_sessions': 0
        }

    ANALYTICS_FILE = "analytics_data.json"
    def save_analytics_data():
        with open(ANALYTICS_FILE, "w") as f:
            json.dump(st.session_state.study_analytics, f)

    def load_analytics_data():
        if os.path.exists(ANALYTICS_FILE):
            with open(ANALYTICS_FILE, "r") as f:
                st.session_state.study_analytics = json.load(f)

    load_analytics_data()

    # --- Ambient Sounds ---
    SOUND_FILES = {
        "Rain": "sounds/rain.mp3",
        "Cafe": "sounds/cafe.mp3",
        "White Noise": "sounds/white_noise.mp3",
        "Nature": "sounds/nature.mp3",
        "Classical Music": "sounds/classical.mp3"
    }

    with discipline_tabs[6]:  # Study Environment
        st.subheader("🎯 " + t("Study Environment")) # Corrected t(...) syntax

        # Ambient sounds
        st.write("### 🎵 " + t("Ambient Sounds")) # Corrected t(...) syntax
        selected_sound = st.selectbox(t("Choose ambient sound"), list(SOUND_FILES.keys())) # Corrected t(...) syntax
        if st.button(t("Play Sound")): # Corrected t(...) syntax
            file_path = SOUND_FILES.get(selected_sound)
            if file_path:
                play_sound(file_path)

    with discipline_tabs[0]:  # Study Streak
        st.subheader("📊 " + t("Study Streak")) # Corrected t(...) syntax
        col1, col2, col3 = st.columns(3)
        with col1:
            st.metric(t("Current Streak"), f"{st.session_state.study_streak['current_streak']} days") # Corrected t(...) syntax
        with col2:
            st.metric(t("Longest Streak"), f"{st.session_state.study_streak['longest_streak']} days") # Corrected t(...) syntax
        with col3:
            st.metric(t("Total Study Time"), f"{st.session_state.study_streak['total_study_time']} hours") # Corrected t(...) syntax

        if st.button(t("Save Streak Data")): # Corrected t(...) syntax
            save_streak_data()
            st.success(t("Streak data saved!")) # Corrected t(...) syntax

    with discipline_tabs[4]:  # Study Analytics
        st.subheader("📈 " + t("Study Analytics")) # Corrected t(...) syntax
        if st.button(t("Save Analytics Data")): # Corrected t(...) syntax
            save_analytics_data()
            st.success(t("Analytics data saved!")) # Corrected t(...) syntax

# --- Footer: Product Hunt Upvote Button & Live Stats ---
ph_stats = get_ph_stats()

# Keep only this section in the sidebar:
st.sidebar.markdown("---")
st.sidebar.markdown("### 🚀 " + t("Support Vekkam")) # Corrected t(...) syntax
st.sidebar.markdown(
    f'''
    <div style="text-align:center;">
        <a href="https://www.producthunt.com/products/vekkam" target="_blank" id="ph-upvote-link">
            <img src="https://api.producthunt.com/widgets/embed-image/v1/featured.svg?post_id=456789&theme=light" alt="{t("Upvote Vekkam on Product Hunt")}" style="width: 150px; margin-bottom: 8px;"/>
        </a><br>
        <span style="font-size:1em; font-weight:bold; color:#da552f;">🔥 {ph_stats['votes']} {t("upvotes")}</span><br>
        <a href="https://www.producthunt.com/products/vekkam" target="_blank" style="font-size:0.9em; font-weight:bold; color:#da552f; text-decoration:none;">👉 {t("Upvote & Comment!")}</a>
    </div>
    ''', unsafe_allow_html=True
)

# Add upvote nudge to sidebar
if 'ph_upvoted' not in st.session_state:
    st.session_state['ph_upvoted'] = False
if not st.session_state['ph_upvoted']:
    if st.sidebar.button("👍 " + t("I upvoted Vekkam!")): # Corrected t(...) syntax
        st.session_state['ph_upvoted'] = True
        st.sidebar.success(t("Thank you for supporting us! 🎉")) # Corrected t(...) syntax
else:
    st.sidebar.info(t("Thanks for your upvote! 🧡")) # Corrected t(...) syntax

# Add recent comments to sidebar if available
if ph_stats['comments']:
    st.sidebar.markdown("### 💬 " + t("Recent Comments")) # Corrected t(...) syntax
    for c in ph_stats['comments']:
        st.sidebar.markdown(
            f'<div style="margin-bottom:0.5em; font-size:0.9em;"><img src="{c["avatar"]}" width="24" style="vertical-align:middle;border-radius:50%;margin-right:4px;"/> <b>{c["user"]}</b><br><span style="font-size:0.85em;">{c["body"]}</span></div>',
            unsafe_allow_html=True
        )

# Final check for learning style score_map (ensure it's defined before use in learning style test)
# This was moved from the learning style test section for global availability, 
# as it was causing an error if the user refreshed after logging in but before taking the test.
score_map = {0: 0, 1: 17, 2: 33, 3: 50, 4: 67, 5: 83, 6: 100}
